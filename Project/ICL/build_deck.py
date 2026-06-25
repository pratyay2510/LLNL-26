# -*- coding: utf-8 -*-
"""Build the ICL literature-survey deck (light mode)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette (light mode) ----------
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x20, 0x24, 0x2B)   # near-black body
NAVY    = RGBColor(0x1A, 0x2B, 0x4A)   # titles
BLUE    = RGBColor(0x1A, 0x73, 0xE8)   # accent (google blue)
TEAL    = RGBColor(0x12, 0x9D, 0x8E)
AMBER   = RGBColor(0xE8, 0x71, 0x0A)
GREY    = RGBColor(0x5F, 0x66, 0x70)
LIGHT   = RGBColor(0xF1, 0xF4, 0xF9)   # light panel
LIGHTER = RGBColor(0xF7, 0xF9, 0xFC)
CALLBG  = RGBColor(0xE8, 0xF0, 0xFE)   # callout bg (light blue)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s

def _set(run, size, color, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font

def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf

def rect(s, l, t, w, h, fill, line=None, round=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def add_runs(p, segments, size, base_color):
    """segments: list of (text, bold, italic) tuples."""
    for seg in segments:
        txt, bold = seg[0], seg[1]
        italic = seg[2] if len(seg) > 2 else False
        r = p.add_run(); r.text = txt
        _set(r, size, BLUE if bold else base_color, bold=bold, italic=italic)

def parse_bold(text):
    """split on **bold** and *italic* into (text,bold,italic) segments."""
    out = []
    for chunk, bold in [(c, i % 2 == 1) for i, c in enumerate(text.split("**"))]:
        if chunk == "":
            continue
        if bold:
            out.append((chunk, True, False))
        else:
            for it_chunk, italic in [(c, j % 2 == 1)
                                     for j, c in enumerate(chunk.split("*"))]:
                if it_chunk == "":
                    continue
                out.append((it_chunk, False, italic))
    return out

# ---------- header strip ----------
def header(s, kicker, title, accent=BLUE):
    rect(s, 0, 0, Inches(0.22), SH, accent)            # left spine
    _, tf = textbox(s, Inches(0.6), Inches(0.34), Inches(12.2), Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = kicker.upper()
    _set(r, 12.5, accent, bold=True)
    r.font.name = "Calibri"
    _, tf2 = textbox(s, Inches(0.6), Inches(0.72), Inches(12.2), Inches(0.95))
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = title
    _set(r2, 27, NAVY, bold=True)
    rect(s, Inches(0.62), Inches(1.62), Inches(1.5), Pt(3), accent)

# ---------- bullets ----------
def bullets(s, items, left=Inches(0.62), top=Inches(1.95),
            width=Inches(7.6), size=16, gap=10):
    """items: list of (level, text) ; text supports **bold**. level 0/1."""
    tb, tf = textbox(s, left, top, width, SH - top - Inches(0.4))
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if level == 0 else gap - 4)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        if level == 0:
            bullet = "▶  "; bcol = BLUE; bsize = size; tcol = INK
        else:
            bullet = "     –  "; bcol = GREY; bsize = size - 2; tcol = GREY
        rb = p.add_run(); rb.text = bullet
        _set(rb, bsize, bcol, bold=(level == 0))
        add_runs(p, parse_bold(text), bsize, tcol)
    return tf

# ---------- callout (key takeaway) ----------
def callout(s, text, top=None, left=Inches(0.62), width=Inches(12.1),
            bg=CALLBG, bar=BLUE):
    if top is None:
        top = SH - Inches(1.25)
    h = Inches(0.85)
    box = rect(s, left, top, width, h, bg, round=True)
    rect(s, left, top, Inches(0.09), h, bar)
    _, tf = textbox(s, left + Inches(0.32), top, width - Inches(0.5), h,
                    anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Takeaway   "
    _set(r, 13, bar, bold=True)
    add_runs(p, parse_bold(text), 14, INK)

# ---------- simple horizontal pipeline diagram ----------
def pipeline(s, steps, top, left=Inches(0.62), total_w=Inches(12.1),
             box_color=LIGHT, accent=BLUE, h=Inches(0.95)):
    n = len(steps)
    gap = Inches(0.28)
    bw = Emu(int((total_w - gap * (n - 1)) / n))
    x = left
    for i, (label, sub) in enumerate(steps):
        b = rect(s, x, top, bw, h, box_color, line=accent, round=True)
        _, tf = textbox(s, x + Inches(0.08), top, bw - Inches(0.16), h,
                        anchor=MSO_ANCHOR.MIDDLE)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        _set(r, 12.5, NAVY, bold=True)
        if sub:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = sub
            _set(r2, 9.5, GREY)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    x + bw + Inches(0.02), top + h/2 - Inches(0.10),
                                    gap - Inches(0.04), Inches(0.20))
            ar.fill.solid(); ar.fill.fore_color.rgb = accent
            ar.line.fill.background(); ar.shadow.inherit = False
        x = Emu(int(x) + int(bw) + int(gap))

# =====================================================================
#  TITLE SLIDE
# =====================================================================
s = slide()
rect(s, 0, 0, SW, Inches(2.55), NAVY)
rect(s, 0, Inches(2.55), SW, Inches(0.10), BLUE)
_, tf = textbox(s, Inches(0.9), Inches(0.72), Inches(11.5), Inches(1.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "From In-Context Learning to Steering"
_set(r, 40, WHITE, bold=True)
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "A literature survey toward query-conditional latent control"
_set(r2, 19, RGBColor(0xC9, 0xD7, 0xF0))

_, tf3 = textbox(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(2.0))
lines = [
    ("Three papers, one thread:", True),
    ("①  GPT-3 — learning from examples in the prompt", False),
    ("②  In-Context Vectors — moving the task signal into the activations", False),
    ("③  LIVE — learning that vector for harder, multimodal tasks", False),
    ("…then: open questions toward a possible problem statement", False),
]
first = True
for txt, head in lines:
    p = tf3.paragraphs[0] if first else tf3.add_paragraph()
    first = False
    p.space_after = Pt(7)
    r = p.add_run(); r.text = txt
    _set(r, 16.5 if head else 15, NAVY if head else INK, bold=head)
_, tf4 = textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
r = tf4.paragraphs[0].add_run()
r.text = "Research discussion  ·  In-Context Learning / Activation Steering"
_set(r, 12.5, GREY, italic=True)

# =====================================================================
#  SECTION DIVIDER helper
# =====================================================================
def divider(num, title, subtitle, accent=BLUE):
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, Inches(0.9), Inches(2.6), Inches(1.7), Pt(4), accent)
    _, tf = textbox(s, Inches(0.9), Inches(1.7), Inches(11), Inches(0.8))
    r = tf.paragraphs[0].add_run(); r.text = num
    _set(r, 16, accent, bold=True)
    _, tf2 = textbox(s, Inches(0.9), Inches(2.85), Inches(11.4), Inches(1.6))
    r = tf2.paragraphs[0].add_run(); r.text = title
    _set(r, 34, WHITE, bold=True)
    _, tf3 = textbox(s, Inches(0.9), Inches(4.25), Inches(11.0), Inches(1.2))
    r = tf3.paragraphs[0].add_run(); r.text = subtitle
    _set(r, 18, RGBColor(0xBF, 0xD0, 0xEA))
    return s

# =====================================================================
#  PAPER 1 — GPT-3
# =====================================================================
divider("PAPER 1", "Language Models are Few-Shot Learners",
        "Brown et al., NeurIPS 2020  —  GPT-3  ·  the case for learning in the prompt")

# 1.1 Motivation
s = slide()
header(s, "GPT-3  ·  Motivation", "The fine-tuning bottleneck")
bullets(s, [
    (0, "Every new task needed its **own labeled dataset** and a round of **gradient updates**."),
    (1, "Costly to collect data; a fresh fine-tuned model per task."),
    (1, "Narrow fine-tuning can **overfit** and generalize poorly out-of-distribution."),
    (0, "**Humans don't work this way** — a short instruction or a couple of examples is enough."),
    (0, "The goal: **one task-agnostic model that adapts at inference**, with no weight updates."),
], width=Inches(7.4))
# side visual
rect(s, Inches(8.35), Inches(2.05), Inches(4.35), Inches(3.7), LIGHTER, line=LIGHT, round=True)
_, tf = textbox(s, Inches(8.6), Inches(2.25), Inches(3.9), Inches(0.5))
r = tf.paragraphs[0].add_run(); r.text = "Two ways to teach a task"
_set(r, 13, NAVY, bold=True)
rect(s, Inches(8.6), Inches(2.8), Inches(3.85), Inches(1.15), RGBColor(0xFB,0xE9,0xE7), line=AMBER, round=True)
_, tf = textbox(s, Inches(8.75), Inches(2.9), Inches(3.6), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; r=p.add_run(); r.text="Fine-tuning"; _set(r,12.5,AMBER,bold=True)
p=tf.add_paragraph(); r=p.add_run(); r.text="new data + gradient updates, per task"; _set(r,10.5,GREY)
rect(s, Inches(8.6), Inches(4.1), Inches(3.85), Inches(1.45), CALLBG, line=BLUE, round=True)
_, tf = textbox(s, Inches(8.75), Inches(4.2), Inches(3.6), Inches(1.25), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; r=p.add_run(); r.text="In-context learning"; _set(r,12.5,BLUE,bold=True)
p=tf.add_paragraph(); r=p.add_run(); r.text="describe the task + show a few examples in the prompt — no updates"; _set(r,10.5,GREY)
callout(s, "Reframe adaptation from **changing weights** to **conditioning on a prompt**.")

# 1.2 Method
s = slide()
header(s, "GPT-3  ·  Method", "Scale + a prompt, instead of training")
bullets(s, [
    (0, "**Scale is the lever:** a 175B-parameter autoregressive LM, ~10× larger than prior models."),
    (0, "Adapt by **conditioning on a prompt** — a task description and/or K worked examples."),
    (1, "Demonstrations sit in the context window; the model infers the pattern and completes the query."),
    (0, "Four regimes compared — **only fine-tuning touches the weights:**"),
], top=Inches(1.9), width=Inches(12.0))
# four regimes row
pipeline(s, [
    ("Fine-tuning", "gradient updates"),
    ("Few-shot", "K examples in prompt"),
    ("One-shot", "1 example"),
    ("Zero-shot", "instruction only"),
], top=Inches(4.0))
callout(s, "Same frozen model, same architecture — **the prompt selects the task**.",
        top=Inches(5.35))

# 1.3 Method
s = slide()
header(s, "GPT-3  ·  Method", "Why it works: ICL as fast 'meta-learning'")
bullets(s, [
    (0, "**Two learning loops:** slow pre-training (outer loop) builds broad skills; **fast in-context adaptation** (inner loop) picks the task at runtime."),
    (0, "Demonstrations steer the model **implicitly** — attention shifts the internal activations toward the demonstrated task."),
    (0, "In-context learning is an **emergent** ability: it sharpens as models get bigger."),
    (1, "The gap between zero-, one-, and few-shot **widens with scale** — larger models are better in-context learners."),
], width=Inches(12.0))
callout(s, "The 'learning' happens **inside the activations at inference** — no parameters move.")

# 1.4 Method / results
s = slide()
header(s, "GPT-3  ·  Results", "One model, many tasks — from text alone")
bullets(s, [
    (0, "Few-shot accuracy **rises with both model size and the number of demonstrations**."),
    (0, "Competitive with — occasionally beating — **fine-tuned SOTA** (e.g. closed-book TriviaQA)."),
    (0, "**No task-specific architecture or data pipeline** — tasks are specified purely in text."),
    (1, "Translation, QA, cloze, comprehension — all from the same weights."),
], width=Inches(12.0))
callout(s, "Proof of concept: **a general model adapts on the fly** — the prompt is the interface.")

# 1.5 Shortcomings
s = slide()
header(s, "GPT-3  ·  Shortcomings", "What in-context learning pays for", accent=AMBER)
bullets(s, [
    (0, "**Demonstrations cost context tokens** — the window is bounded (2048), so more examples = more compute and a hard ceiling."),
    (0, "**Brittle and uneven** — performance swings with the choice of template, verbalizer, and the selection & order of examples."),
    (0, "**Control is implicit** — you cannot dial the strength of the task signal; attention decides how much to shift."),
], width=Inches(12.0), size=17, gap=14)
callout(s, "The task signal is **real but trapped in the prompt** — expensive, fragile, and not tunable.",
        bg=RGBColor(0xFD,0xEF,0xE3), bar=AMBER)

# 1.6 Bridge
s = slide()
header(s, "GPT-3 → next", "If the signal lives in the activations…", accent=AMBER)
bullets(s, [
    (0, "GPT-3 shows the task signal is **carried by the model's internal states**, but delivers it the slow way — through tokens in the prompt."),
    (0, "**Open question:** can we **extract that signal once** and inject it directly into the activations — no demos in the context, and with a strength we control?"),
], width=Inches(12.0), size=17, gap=14)
callout(s, "Next: **In-Context Vectors** — turn the prompt's effect into a single steerable direction.",
        bg=CALLBG, bar=BLUE)

# =====================================================================
#  PAPER 2 — ICV
# =====================================================================
divider("PAPER 2", "In-Context Vectors (ICV)",
        "Liu et al., ICML 2024  —  making ICL effective & controllable via latent-space steering")

# 2.1 Motivation
s = slide()
header(s, "ICV  ·  Motivation", "Make the hidden 'shift' explicit")
bullets(s, [
    (0, "ICL is effective but **token-hungry, brittle, and hard to control** (the GPT-3 pain points)."),
    (0, "**Key insight:** in-context learning is really just **'shifting' the latent states** — attention nudges activations from the input toward the demonstrated answer."),
    (0, "Today that shift is **silent and automatic**. So: **can we make it explicit, controllable, and free of prompt tokens?**"),
], width=Inches(12.0), size=17, gap=14)
callout(s, "Stop paying for demonstrations every query — **summarize them once into a vector**.")

# 2.2 Method
s = slide()
header(s, "ICV  ·  Method", "From demonstrations to one direction")
bullets(s, [
    (0, "Recast ICL as a single **In-Context Vector (ICV)** in latent space."),
    (0, "**Step 1 — Task summary:** one forward pass over the demos; read the latent states for each input x and target y."),
    (0, "ICV = the **first principal direction (PCA) of h(y) − h(x)** — the dominant 'input → target' move."),
    (1, "Training-free: no fine-tuning, no extra parameters — just a PCA over a handful of demos."),
], width=Inches(12.0))
pipeline(s, [
    ("Demos (x, y)", "few examples"),
    ("Forward pass", "read latent states"),
    ("h(y) − h(x)", "per-pair difference"),
    ("PCA → ICV", "principal direction"),
], top=Inches(4.55))
callout(s, "The whole demonstration set collapses into **one direction in activation space**.",
        top=Inches(5.7))

# 2.3 Method
s = slide()
header(s, "ICV  ·  Method", "Steer the query, with a knob")
bullets(s, [
    (0, "**Step 2 — Feature shifting:** add the ICV to the query's latent states across the model, then generate."),
    (0, "The prompt holds **no demonstrations** — just the query. Demos already live in the vector."),
    (0, "**Strength is a dial:** scale the ICV's magnitude to turn task adherence up or down."),
    (1, "Bounded context is no longer a limit — a vector can summarize **more demos than would fit**."),
], width=Inches(12.0))
callout(s, "Explicit, tunable control — **you set how hard to steer**, instead of leaving it to attention.")

# 2.4 Method / results
s = slide()
header(s, "ICV  ·  Results", "Cheaper, stronger, composable")
bullets(s, [
    (0, "**Near-zero overhead:** demos summarized once; no per-query token cost."),
    (0, "**Beats ICL and LoRA** on safety/detox, style transfer, role-play, and formatting (Falcon, Llama)."),
    (0, "**Skills compose by vector arithmetic** — add ICVs to follow several instructions at once."),
], width=Inches(12.0), size=17, gap=13)
pipeline(s, [
    ("Detox vector", ""),
    ("+  Formality vector", ""),
    ("=  Combined steer", "one query, two skills"),
], top=Inches(4.6), box_color=LIGHTER)
callout(s, "Latent steering = **ICL's effect, minus the prompt cost, plus an explicit control surface**.",
        top=Inches(5.75))

# 2.5 Shortcomings
s = slide()
header(s, "ICV  ·  Shortcomings", "One global vector for everything", accent=AMBER)
bullets(s, [
    (0, "**A single vector per task** — the *same* shift is applied to *every* query, regardless of what it is asking."),
    (0, "The **non-learnable PCA summary** is a coarse average; it can't capture a task with many internal sub-types."),
    (0, "**Works on 'easy', uniform patterns** (antonyms, country→capital) but frays when one task hides many question types."),
], width=Inches(12.0), size=17, gap=14)
callout(s, "A global average **helps some inputs and silently hurts others** — and PCA can't carry rich tasks.",
        bg=RGBColor(0xFD,0xEF,0xE3), bar=AMBER)

# 2.6 Bridge
s = slide()
header(s, "ICV → next", "When one direction isn't enough", accent=AMBER)
bullets(s, [
    (0, "A fixed PCA direction **can't represent a task that is really many tasks** — e.g. Visual QA, where 'what is this?' and 'how many?' need different abilities."),
    (0, "**Open question:** instead of *reading off* a direction with PCA, can we **learn** a vector that distills the task signal PCA misses?"),
], width=Inches(12.0), size=17, gap=14)
callout(s, "Next: **LIVE** — a learnable in-context vector for harder, multimodal tasks.",
        bg=CALLBG, bar=BLUE)

# =====================================================================
#  PAPER 3 — LIVE
# =====================================================================
divider("PAPER 3", "LIVE: Learnable In-Context Vector for VQA",
        "Peng et al., NeurIPS 2024  —  learning the steering vector for multimodal tasks")

# 3.1 Motivation
s = slide()
header(s, "LIVE  ·  Motivation", "Steering vectors for multimodal models")
bullets(s, [
    (0, "Extend in-context vectors to **Large Multimodal Models** (e.g. IDEFICS/Flamingo), where ICL is **even costlier and more brittle** — images make demonstrations huge and selection harder."),
    (0, "**Non-learnable ICVs (PCA / mean-diff) collapse on VQA:** one vector can't cover a task that bundles classification, counting, reading, and more."),
    (0, "**Open question:** can the vector be **learned** to distill task information a fixed PCA direction simply cannot?"),
], width=Inches(12.0), size=16.5, gap=13)
callout(s, "Same idea — one steering vector — but **trained**, so it can carry a richer task.")

# 3.2 Method
s = slide()
header(s, "LIVE  ·  Method", "Distill demonstrations into a learned vector")
bullets(s, [
    (0, "Train a shift vector to **mimic the effect of demonstrations** — no demos at inference."),
    (0, "**Distillation objective:** match the model's output distribution with *LIVE only* to its output with *32-shot demos*."),
    (0, "**Randomize the 32-shot demos per query** during training → the vector keeps the **task signal** and drops example-specific noise."),
], width=Inches(12.0))
pipeline(s, [
    ("32-shot demos", "randomized per query"),
    ("Match output dist.", "LIVE vs demos"),
    ("Learn vector", "task signal only"),
    ("Inference", "add vector, no demos"),
], top=Inches(4.55))
callout(s, "PCA *reads off* a direction; LIVE **optimizes** one to reproduce the demonstrations' effect.",
        top=Inches(5.7))

# 3.3 Method
s = slide()
header(s, "LIVE  ·  Method", "Per-layer vectors, learned strength")
bullets(s, [
    (0, "**A separate vector per layer** — different layers play different roles, so each gets its own fine-grained shift."),
    (0, "**Learnable magnitude α** — the steering strength is *trained*, not hand-tuned."),
    (0, "Only the vectors are trained; the **backbone stays frozen** — still lightweight."),
], width=Inches(12.0), size=17, gap=14)
callout(s, "More expressive than a single global PCA direction, while keeping the **frozen-model, add-a-vector** recipe.")

# 3.4 Method / results
s = slide()
header(s, "LIVE  ·  Results", "Big efficiency, better accuracy")
bullets(s, [
    (0, "**~25× fewer FLOPs** than 32-shot ICL at matched accuracy (≈ 1/24.97)."),
    (0, "**+2.36 / +1.6 accuracy** on VQAv2 / OKVQA over 32-shot ICL."),
    (0, "**Far more sample-efficient than LoRA** — ~500 vs ~8000 training samples for comparable trainable parameters."),
], width=Inches(12.0), size=17, gap=14)
callout(s, "Learning the vector **rescues steering on hard, heterogeneous tasks** that broke non-learnable ICVs.")

# 3.5 Shortcomings
s = slide()
header(s, "LIVE  ·  Shortcomings", "Still global, still always-on", accent=AMBER)
bullets(s, [
    (0, "**One learned vector per task** — better than PCA, but still **not query-conditional**: every input gets the same shift."),
    (0, "**Needs an offline training loop** with labeled demos; **adding a new skill means retraining**."),
    (0, "**No sense of *when* to steer** — the shift is always applied, even on inputs where it could **flip a correct answer to wrong**."),
], width=Inches(12.0), size=17, gap=14)
callout(s, "We made the vector smarter — but it's still **one-size-fits-all and unconditionally on**.",
        bg=RGBColor(0xFD,0xEF,0xE3), bar=AMBER)

# 3.6 Bridge / synthesis
s = slide()
header(s, "Across all three", "The thread — and what's still missing", accent=TEAL)
# two-column comparison
rect(s, Inches(0.62), Inches(1.95), Inches(6.0), Inches(3.7), LIGHTER, line=LIGHT, round=True)
_, tf = textbox(s, Inches(0.85), Inches(2.1), Inches(5.6), Inches(3.4))
p=tf.paragraphs[0]; r=p.add_run(); r.text="What got better"; _set(r,16,TEAL,bold=True)
for t in ["GPT-3 → task signal lives in the activations",
          "ICV → extract it into one explicit, tunable vector",
          "LIVE → learn that vector for rich, hard tasks",
          "Cost fell; control rose."]:
    p=tf.add_paragraph(); p.space_before=Pt(6)
    r=p.add_run(); r.text="✓  "+t; _set(r,13.5,INK)
rect(s, Inches(6.92), Inches(1.95), Inches(5.8), Inches(3.7), RGBColor(0xFD,0xEF,0xE3), line=AMBER, round=True)
_, tf = textbox(s, Inches(7.15), Inches(2.1), Inches(5.4), Inches(3.4))
p=tf.paragraphs[0]; r=p.add_run(); r.text="What stayed the same"; _set(r,16,AMBER,bold=True)
for t in ["The vector is global — one shift for all queries",
          "Steering is always on — no notion of when",
          "Adding a skill needs (re)training",
          "Steering can silently harm some inputs."]:
    p=tf.add_paragraph(); p.space_before=Pt(6)
    r=p.add_run(); r.text="•  "+t; _set(r,13.5,INK)
callout(s, "What if steering were **retrieved per query, composed, and gated** to avoid harm?",
        bg=RGBColor(0xE7,0xF6,0xF4), bar=TEAL)

# =====================================================================
#  OPEN QUESTIONS SECTION
# =====================================================================
divider("DISCUSSION", "Open Questions & a Possible Direction",
        "Framing a problem statement — avenues to explore, not committed answers", accent=TEAL)

# Open questions
s = slide()
header(s, "Discussion", "Open questions worth posing", accent=TEAL)
bullets(s, [
    (0, "Can the **prompt itself act as a router** — signalling *which* steering vector to use? (a math query pulls a math vector)"),
    (0, "Can we keep a **bank of skill vectors** and **retrieve-then-steer** per query, instead of one global average?"),
    (0, "Can steering be **query-conditional** *without retraining* to add a new skill?"),
    (0, "Can we **steer only when it's safe**, and **abstain** otherwise?"),
    (0, "Can we **bound the harm rate** — the chance steering flips a correct answer to wrong — at a chosen level?"),
    (0, "Is **heterogeneity** (a mix of query types) the regime where a single global vector *must* fail?"),
    (0, "How do we **compose** several skill vectors **without interference**?"),
], width=Inches(12.0), size=15.5, gap=9)
callout(s, "Analogy to test: RAG decouples **knowledge** from weights — can steering decouple **skills** from weights?",
        bg=RGBColor(0xE7,0xF6,0xF4), bar=TEAL)

# Possible direction
s = slide()
header(s, "Discussion", "A possible direction (a sketch, not a commitment)", accent=TEAL)
bullets(s, [
    (0, "**Retrieval-Augmented Steering (RAS)** — frame steering as *routing + abstention*, by analogy to RAG."),
    (1, "**Build** an open-domain bank of skill vectors offline (cluster demos → one vector per skill)."),
    (1, "**Retrieve** the nearest skill vector(s) for each query → **compose** → **steer-or-abstain**."),
    (0, "**Near-zero cost:** one query embedding + a lookup + a vector add. New skill = new bank entry, **no retraining**."),
    (0, "Headline question to probe: a **risk–coverage frontier** — *how much gain at a guaranteed harm rate?*"),
], width=Inches(12.0), size=16, gap=11)
pipeline(s, [
    ("Query", "embed"),
    ("Retrieve", "nearest skill vec(s)"),
    ("Compose", "top-k, small"),
    ("Gate", "steer or abstain"),
    ("Steer", "add vector"),
], top=Inches(5.25), accent=TEAL)
# note line
_, tf = textbox(s, Inches(0.62), Inches(6.45), Inches(12), Inches(0.5))
r=tf.paragraphs[0].add_run()
r.text="Posed as open avenues — start from the simplest defensible choice in each block and ablate upward."
_set(r,12.5,GREY,italic=True)

# Reading list
s = slide()
header(s, "Discussion", "Reading list the direction builds on", accent=TEAL)
def readcol(left, title, items, col):
    rect(s, left, Inches(1.95), Inches(3.85), Inches(4.55), LIGHTER, line=LIGHT, round=True)
    _, tf = textbox(s, left+Inches(0.22), Inches(2.1), Inches(3.45), Inches(4.3))
    p=tf.paragraphs[0]; r=p.add_run(); r.text=title; _set(r,13.5,col,bold=True)
    for it in items:
        p=tf.add_paragraph(); p.space_before=Pt(5); p.line_spacing=1.0
        r=p.add_run(); r.text="•  "+it; _set(r,11,INK)
readcol(Inches(0.62), "Closest works (conditional steering)", [
    "K-CAST — kNN-conditional steering",
    "CAST — conditional activation steering",
    "HyperSteer — generate vectors via hypernet",
    "Steer2Adapt — compose a fixed basis",
    "PDS — prototype-based dynamic steering",
], BLUE)
readcol(Inches(4.74), "Primitives we build on", [
    "CAA — mean-difference steering vector",
    "ICV — PCA of h(y) − h(x)  (Paper 2)",
    "LIVE — learned ICV for VQA  (Paper 3)",
    "SADI — semantics-adaptive intervention",
    "Composition — MAT-Steer, K-Steering",
], TEAL)
readcol(Inches(8.86), "No-harm / selective prediction", [
    "Conformal Abstention for LLMs",
    "Conformal Risk Control for LLMs",
    "'Entropy alone is insufficient' for gating",
    "Steering Externalities — benign steering risk",
    "Positioning: RAG-analogs — LoRA-Retriever, EPR",
], AMBER)
_, tf = textbox(s, Inches(0.62), Inches(6.65), Inches(12), Inches(0.5))
r=tf.paragraphs[0].add_run()
r.text="Verify each exact reference/venue before depending on it."
_set(r,12,GREY,italic=True)

# Closing
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, Inches(0.9), Inches(2.9), Inches(1.7), Pt(4), TEAL)
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.4), Inches(1.0))
r=tf.paragraphs[0].add_run(); r.text="Where we'd like the discussion to go"
_set(r,28,WHITE,bold=True)
_, tf = textbox(s, Inches(0.9), Inches(3.25), Inches(11.0), Inches(3.0))
for t in ["The task signal moved from the prompt → into one vector → into a learned vector.",
          "It is still global and always-on. The open wedge: retrieve, compose, and gate it.",
          "Goal of today: pressure-test whether 'retrieve-then-steer, safely' is a problem worth owning."]:
    p=tf.paragraphs[0] if t.startswith("The task") else tf.add_paragraph()
    p.space_after=Pt(12)
    r=p.add_run(); r.text="—  "+t; _set(r,16.5,RGBColor(0xCF,0xDD,0xF2))

prs.save("ICL_Literature_Survey.pptx")
print("saved ICL_Literature_Survey.pptx ; slides =", len(prs.slides._sldIdLst))
