"""
Mentor Presentation: ICL Research Plan
Section 1: QVM-FRE-BI Framework for ICL
Section 2: LLMs Are Invertible -- Connection to ICL
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SW, SH = Inches(13.333), Inches(7.5)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x16, 0x21, 0x3E)
BLUE = RGBColor(0x00, 0x8C, 0xFF)
CYAN = RGBColor(0x00, 0xD4, 0xAA)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xAE, 0xC0)
MG = RGBColor(0x60, 0x70, 0x88)

def bg(s, c=DARK):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.color.rgb = c

def bx(s, l, t, w, h, f=CARD, b=None, bw=Pt(1)):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = f
    r.line.color.rgb = b if b else f
    if b: r.line.width = bw

def tx(s, l, t, w, h, txt, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = "Calibri"; p.alignment = a

def ml(s, l, t, w, h, lines):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln[0]; p.font.size = Pt(ln[1])
        p.font.color.rgb = ln[2]; p.font.bold = ln[3] if len(ln)>3 else False
        p.font.name = "Calibri"; p.alignment = ln[4] if len(ln)>4 else PP_ALIGN.LEFT
        p.space_after = Pt(2)

def hdr(s, title, color):
    bx(s, 0, 0, SW, Pt(5), f=color)
    tx(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6), title, 32, WHITE, True)
    bx(s, Inches(0.8), Inches(1.0), Inches(2.2), Pt(3), f=color)

def main():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    L6 = prs.slide_layouts[6]

    # ===== SLIDE 1: TITLE =====
    s = prs.slides.add_slide(L6); bg(s)
    bx(s, 0, 0, SW, Pt(5), f=BLUE)
    tx(s, Inches(1.2), Inches(1.4), Inches(11), Inches(0.7),
       "In-Context Learning as", 38, GRAY)
    tx(s, Inches(1.2), Inches(2.1), Inches(11), Inches(0.9),
       "Latent Space Manipulation", 48, WHITE, True)
    bx(s, Inches(1.2), Inches(3.2), Inches(2.5), Pt(4), f=CYAN)
    tx(s, Inches(1.2), Inches(3.5), Inches(11), Inches(0.5),
       "Bridging Classical RF Theory with Modern Transformer Mechanics", 19, CYAN)
    ml(s, Inches(1.2), Inches(4.6), Inches(11), Inches(2.0), [
        ("Pratyay Dutta", 22, WHITE, True),
        ("University of California, Riverside  |  Lawrence Livermore National Laboratory", 15, GRAY),
        ("", 8, GRAY),
        ("Advisor: Prof. Bir Bhanu  |  June 2026", 16, BLUE),
    ])

    # ===== SLIDE 2: SECTION DIVIDER 1 =====
    s = prs.slides.add_slide(L6); bg(s)
    bx(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5), f=CARD, b=BLUE, bw=Pt(2))
    tx(s, Inches(1), Inches(2.7), Inches(11), Inches(0.6),
       "Section 1", 20, BLUE)
    tx(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
       "The QVM-FRE-BI Framework for ICL", 38, WHITE, True)
    tx(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
       "Connecting classical relevance feedback to modern transformer attention", 16, GRAY)

    # ===== SLIDE 3: THE ANALOGY =====
    s = prs.slides.add_slide(L6); bg(s)
    hdr(s, "Core Insight: ICL ~ Relevance Feedback", GREEN)

    # QVM box
    bx(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.5), b=BLUE, bw=Pt(2), f=CARD)
    ml(s, Inches(0.7), Inches(1.5), Inches(3.4), Inches(5.3), [
        ("QVM", 26, BLUE, True, PP_ALIGN.CENTER),
        ("Query Vector Movement", 12, GRAY, False, PP_ALIGN.CENTER),
        ("", 6, GRAY),
        ("Classical (Rocchio):", 13, CYAN, True),
        ("q_new = a*q_old", 13, WHITE),
        ("     + b * mean(D+)", 13, WHITE),
        ("     - g * mean(D-)", 13, WHITE),
        ("", 6, GRAY),
        ("ICL Attention:", 13, CYAN, True),
        ("h_new = h_old", 13, WHITE),
        ("  + Sum attn(QK/sqrt(d)) * V", 13, WHITE),
        ("", 6, GRAY),
        ("Residual = a*q_old       [OK]", 12, GREEN),
        ("Attn weights = b*rel     [OK]", 12, GREEN),
        ("No gamma repulsion      [GAP]", 12, RED),
    ])

    # FRE box
    bx(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), b=PURPLE, bw=Pt(2), f=CARD)
    ml(s, Inches(4.9), Inches(1.5), Inches(3.4), Inches(5.3), [
        ("FRE", 26, PURPLE, True, PP_ALIGN.CENTER),
        ("Feature Relevance Estimation", 12, GRAY, False, PP_ALIGN.CENTER),
        ("", 6, GRAY),
        ("Classical (TPAMI 2005):", 13, CYAN, True),
        ("w_j = (1/sigma_j)", 13, WHITE),
        ("      / Sum_k(1/sigma_k)", 13, WHITE),
        ("", 4, GRAY),
        ("Tight cluster -> high weight", 12, GREEN),
        ("Wide spread -> low weight", 12, ORANGE),
        ("", 6, GRAY),
        ("In Transformers:", 13, CYAN, True),
        ("Multi-head attn = subspace FRE", 12, WHITE),
        ("LayerNorm = axis rescaling", 12, WHITE),
        ("Implicit but uncontrolled  [??]", 12, ORANGE),
    ])

    # BI box
    bx(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.5), b=CYAN, bw=Pt(2), f=CARD)
    ml(s, Inches(9.1), Inches(1.5), Inches(3.4), Inches(5.3), [
        ("BI", 26, CYAN, True, PP_ALIGN.CENTER),
        ("Bayesian Inference", 12, GRAY, False, PP_ALIGN.CENTER),
        ("", 6, GRAY),
        ("Classical:", 13, CYAN, True),
        ("P(rel|x) ~ P(x|rel)*P(rel)", 13, WHITE),
        ("", 6, GRAY),
        ("ICL:", 13, CYAN, True),
        ("P(y|x,C) =", 13, WHITE),
        ("  integral P(y|x,g)*P(g|C) dg", 13, WHITE),
        ("", 4, GRAY),
        ("Context C infers latent task g", 12, WHITE),
        ("", 6, GRAY),
        ("Well-established analogy  [OK]", 12, GREEN),
        ("Pretraining = prior P(g)  [OK]", 12, GREEN),
    ])

    # ===== SLIDE 4: RESEARCH GAPS =====
    s = prs.slides.add_slide(L6); bg(s)
    hdr(s, "Identified Gaps in ICL via CBIR Lens", RED)

    gaps = [
        ("Missing Gamma-Term", RED, "ICL has NO negative demo repulsion mechanism", "The g*mean(D-) from Rocchio is absent in attention"),
        ("Implicit FRE Uncontrolled", PURPLE, "Transformers may do FRE via attention heads", "Never verified; no classical 1/sigma_j connection made"),
        ("Positional Bias Corrupts QVM", BLUE, "Attention conflates position with semantic relevance", "Recency bias = position-contaminated softmax weights"),
        ("No Unified Framework", CYAN, "QVM, FRE, BI studied in isolation for ICL", "Our contribution: first unified lens for ICL analysis"),
    ]
    for i, (t, c, l, r) in enumerate(gaps):
        y = Inches(1.3) + i * Inches(1.45)
        bx(s, Inches(0.8), y, Inches(11.5), Inches(1.3), b=c, f=CARD)
        tx(s, Inches(1.0), y + Inches(0.05), Inches(5), Inches(0.3), t, 15, c, True)
        tx(s, Inches(1.0), y + Inches(0.4), Inches(5.3), Inches(0.8), l, 12, GRAY)
        tx(s, Inches(6.8), y + Inches(0.4), Inches(5.2), Inches(0.8), r, 12, WHITE)

    # ===== SLIDE 5: PHASE 1 EXPERIMENTS =====
    s = prs.slides.add_slide(L6); bg(s)
    hdr(s, "Phase 1: Diagnostic Experiments", CYAN)
    tx(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
       "6 modules to expose ICL failure modes through the QVM-FRE lens", 14, GRAY)

    mods = [
        ("M1", "Noise Injection", BLUE, "Gaussian noise in\nvalue vectors"),
        ("M2", "Attn Skewing", CYAN, "Permute demo order\n-> recency bias"),
        ("M3", "Label Flipping", GREEN, "Format vs semantics\nstrength test"),
        ("M4", "FRE Probing", PURPLE, "Per-dim ablation\n-> FRE quality"),
        ("M5", "Gamma-Term", RED, "Negative demos +\nsynthetic repulsion"),
        ("M6", "Task Vectors", ORANGE, "Extract & analyze\ntask vector geometry"),
    ]
    for i, (lbl, t, c, d) in enumerate(mods):
        x = Inches(0.4) + i * Inches(2.1)
        bx(s, x, Inches(1.6), Inches(1.95), Inches(2.8), b=c, bw=Pt(2), f=CARD)
        bx(s, x+Inches(0.08), Inches(1.68), Inches(0.5), Inches(0.3), f=c)
        tx(s, x+Inches(0.08), Inches(1.68), Inches(0.5), Inches(0.3), lbl, 12, WHITE, True, PP_ALIGN.CENTER)
        tx(s, x+Inches(0.08), Inches(2.1), Inches(1.7), Inches(0.5), t, 14, WHITE, True)
        tx(s, x+Inches(0.08), Inches(2.7), Inches(1.7), Inches(1.2), d, 12, GRAY)

    # Bottom: convergence
    bx(s, Inches(0.4), Inches(4.7), Inches(12.4), Inches(2.4), b=MG, f=CARD)
    ml(s, Inches(0.8), Inches(4.85), Inches(11.6), Inches(2.1), [
        ("Key Hypotheses", 17, CYAN, True),
        ("", 3, GRAY),
        ("H1: Negative demos cause query vector to move AWAY from incorrect class centroids", 13, WHITE),
        ("H2: Optimal gamma* exists that maximizes ICL accuracy (task-dependent)", 13, WHITE),
        ("H3: Gamma-term improvement is largest for ambiguous examples near decision boundaries", 13, WHITE),
        ("", 3, GRAY),
        ("Convergence -> Cross-module analysis -> Unified sensitivity matrix -> Problem formulation", 14, CYAN, True),
    ])

    # ===== SLIDE 6: SECTION DIVIDER 2 =====
    s = prs.slides.add_slide(L6); bg(s)
    bx(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5), f=CARD, b=PURPLE, bw=Pt(2))
    tx(s, Inches(1), Inches(2.7), Inches(11), Inches(0.6),
       "Section 2", 20, PURPLE)
    tx(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
       "LLMs Are Invertible -- What It Means for ICL", 36, WHITE, True)
    tx(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
       "Nikolaou, Mencattini et al. (ICLR 2026) -- Injectivity as a tool for our research", 15, GRAY)

    # ===== SLIDE 7: PAPER SUMMARY =====
    s = prs.slides.add_slide(L6); bg(s)
    hdr(s, "Key Result: Transformers Are Injective", PURPLE)

    bx(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.8), b=PURPLE, bw=Pt(2), f=CARD)
    ml(s, Inches(0.8), Inches(1.4), Inches(5.2), Inches(2.6), [
        ("Core Theorem", 18, PURPLE, True),
        ("", 4, GRAY),
        ("s != s'  =>  r(s; theta) != r(s'; theta)", 16, WHITE, True),
        ("", 4, GRAY),
        ("Different prompts -> different last-token", 13, GRAY),
        ("representations, with probability 1", 13, GRAY),
        ("", 4, GRAY),
        ("Holds at initialization AND preserved", 13, CYAN),
        ("through training (any finite T GD steps)", 13, CYAN),
    ])

    bx(s, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.8), b=BLUE, bw=Pt(2), f=CARD)
    ml(s, Inches(7.1), Inches(1.4), Inches(5.2), Inches(2.6), [
        ("Why It Holds", 18, BLUE, True),
        ("", 4, GRAY),
        ("1. Transformers are real-analytic functions", 13, WHITE),
        ("   (composition of analytic components)", 13, GRAY),
        ("", 3, GRAY),
        ("2. Collision set has Lebesgue measure zero", 13, WHITE),
        ("   h(theta) = ||r(s;theta) - r(s';theta)||^2", 13, GRAY),
        ("", 3, GRAY),
        ("3. GD preserves absolute continuity", 13, WHITE),
        ("   -> can't move params into collision set", 13, GRAY),
    ])

    # SIPIT box
    bx(s, Inches(0.5), Inches(4.4), Inches(12.1), Inches(2.7), b=CYAN, bw=Pt(2), f=RGBColor(0x12, 0x1A, 0x35))
    ml(s, Inches(0.8), Inches(4.55), Inches(11.5), Inches(2.4), [
        ("SIPIT: Constructive Inversion", 18, CYAN, True),
        ("", 4, GRAY),
        ("Sequential Inverse Prompt via Iterative Updates -- recovers exact input from hidden states", 14, WHITE),
        ("", 3, GRAY),
        ("Exploits causal structure: h_t depends only on prefix <s1,...,s_{t-1}> and current token s_t", 13, GRAY),
        ("At each step: cycle through vocab candidates until unique match found", 13, GRAY),
        ("100% accuracy on GPT-2, Mistral-7B, LLaMA-3.1-8B  |  <0.22% vocab explored on average", 13, GREEN),
        ("", 3, GRAY),
        ("Implication: Hidden states are NOT abstractions -- they ARE the prompt in disguise", 14, ORANGE, True),
    ])

    # ===== SLIDE 8: CONNECTION TO OUR WORK =====
    s = prs.slides.add_slide(L6); bg(s)
    hdr(s, "What This Means for Our ICL Research", GREEN)

    # Left: theoretical implications
    bx(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), b=GREEN, bw=Pt(2), f=CARD)
    ml(s, Inches(0.8), Inches(1.4), Inches(5.2), Inches(5.3), [
        ("Validates Our Framework", 18, GREEN, True),
        ("", 6, GRAY),
        ("1. Injectivity = Lossless QVM", 15, CYAN, True),
        ("   If r(s) preserves ALL input info, then", 13, GRAY),
        ("   attention updates are information-preserving", 13, GRAY),
        ("   query movements -- exactly what QVM assumes", 13, GRAY),
        ("", 5, GRAY),
        ("2. Strengthens FRE Probing (M4)", 15, CYAN, True),
        ("   Per-dimension ablation is valid because", 13, GRAY),
        ("   each dimension carries unique, non-redundant", 13, GRAY),
        ("   information (no collapse/aliasing)", 13, GRAY),
        ("", 5, GRAY),
        ("3. Task Vectors Are Meaningful (M6)", 15, CYAN, True),
        ("   Injectivity guarantees task vectors in", 13, GRAY),
        ("   hidden states faithfully encode the full", 13, GRAY),
        ("   context -- probing them is well-founded", 13, GRAY),
    ])

    # Right: practical applications
    bx(s, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), b=PURPLE, bw=Pt(2), f=CARD)
    ml(s, Inches(7.1), Inches(1.4), Inches(5.2), Inches(5.3), [
        ("Actionable Ideas", 18, PURPLE, True),
        ("", 6, GRAY),
        ("A. SIPIT as Ground-Truth Oracle", 15, ORANGE, True),
        ("   Use SIPIT to verify that our latent-space", 13, GRAY),
        ("   interventions (gamma injection, FRE weights)", 13, GRAY),
        ("   don't destroy input fidelity", 13, GRAY),
        ("", 5, GRAY),
        ("B. Inversion-Based Sensitivity Analysis", 15, ORANGE, True),
        ("   Perturb hidden states (our M1 noise injection)", 13, GRAY),
        ("   then run SIPIT to see WHICH tokens are", 13, GRAY),
        ("   affected -- spatial sensitivity map", 13, GRAY),
        ("", 5, GRAY),
        ("C. Collision Distance as ICL Metric", 15, ORANGE, True),
        ("   Use min pairwise L2 between demo hidden", 13, GRAY),
        ("   states as a measure of demo diversity --", 13, GRAY),
        ("   connects to demo selection quality", 13, GRAY),
    ])

    # ===== SLIDE 9: TENTATIVE PLAN =====
    s = prs.slides.add_slide(L6); bg(s)
    hdr(s, "Tentative Plan: Moving Forward", ORANGE)

    weeks = [
        ("Wk 1-2", "Foundation", "Model hooks setup\nBaseline ICL evals\nActivation extraction", BLUE),
        ("Wk 3-4", "FRE + Gamma", "M4: FRE probing\nM5: Gamma injection\nCore experiments", PURPLE),
        ("Wk 5-6", "Diagnostics", "M1-M3: Perturbations\nM6: Task vectors\nCross-module analysis", CYAN),
        ("Wk 7-8", "Invertibility", "SIPIT validation\nSensitivity maps\nCollision metrics", GREEN),
        ("Wk 9-10", "Synthesis", "Problem formulation\nPhase 2 directions\nWriteup", ORANGE),
    ]
    for i, (wk, t, d, c) in enumerate(weeks):
        x = Inches(0.3) + i * Inches(2.55)
        bx(s, x, Inches(1.3), Inches(2.35), Inches(2.8), b=c, bw=Pt(2), f=CARD)
        tx(s, x+Inches(0.05), Inches(1.4), Inches(2.2), Inches(0.3), wk, 13, c, True, PP_ALIGN.CENTER)
        tx(s, x+Inches(0.05), Inches(1.75), Inches(2.2), Inches(0.35), t, 15, WHITE, True, PP_ALIGN.CENTER)
        tx(s, x+Inches(0.05), Inches(2.2), Inches(2.2), Inches(1.5), d, 12, GRAY, False, PP_ALIGN.CENTER)

    bx(s, Inches(0.3), Inches(4.2), Inches(12.7), Pt(3), f=ORANGE)

    # Decision tree
    bx(s, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.6), b=MG, f=CARD)
    ml(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.4), [
        ("Phase 2 Directions (Evidence-Dependent)", 17, ORANGE, True),
        ("", 4, GRAY),
        ("IF weak FRE        ->  Inject explicit FRE via learned axis weighting", 13, PURPLE),
        ("IF gamma helps     ->  Contrastive attention / repulsive residuals", 13, RED),
        ("IF recency bias    ->  Positional debiasing with Rocchio-inspired weights", 13, BLUE),
        ("IF unstable vectors ->  Task vector regularization / demo-agnostic encoding", 13, CYAN),
        ("IF invertibility useful -> SIPIT-guided latent interventions", 13, GREEN),
    ])

    # ===== SLIDE 10: THANK YOU =====
    s = prs.slides.add_slide(L6); bg(s)
    bx(s, 0, 0, SW, Pt(5), f=BLUE)
    tx(s, Inches(2), Inches(2.0), Inches(9), Inches(1.0),
       "Thank You", 50, WHITE, True, PP_ALIGN.CENTER)
    bx(s, Inches(5.5), Inches(3.1), Inches(2.3), Pt(4), f=CYAN)
    tx(s, Inches(2), Inches(3.4), Inches(9), Inches(0.6),
       "Questions & Discussion", 26, GRAY, False, PP_ALIGN.CENTER)
    ml(s, Inches(2), Inches(4.5), Inches(9), Inches(2.0), [
        ("Pratyay Dutta", 22, WHITE, True, PP_ALIGN.CENTER),
        ("pdutta@ucr.edu", 16, CYAN, False, PP_ALIGN.CENTER),
        ("", 6, GRAY),
        ("UC Riverside  |  Lawrence Livermore National Laboratory", 15, GRAY, False, PP_ALIGN.CENTER),
    ])

    out = r"c:\Work\LLNL\LLNL-26\Project\ICL\ICL_Mentor_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
