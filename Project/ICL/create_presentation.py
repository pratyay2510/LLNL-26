"""
Generate research presentation — Reliable version using only textboxes and simple rectangles.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
DARK = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x16, 0x21, 0x3E)
BLUE = RGBColor(0x00, 0x8C, 0xFF)
CYAN = RGBColor(0x00, 0xD4, 0xAA)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xAE, 0xC0)
MGRAY = RGBColor(0x60, 0x70, 0x88)


def bg(slide, color=DARK):
    # Draw a full-slide rectangle covering the entire slide to act as the background.
    # This is highly compatible with all PowerPoint viewers and prevents blank slide issues.
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.color.rgb = color
    return s


def rect(slide, l, t, w, h, fill=CARD, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.color.rgb = fill
    return s


def tb(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def mtb(slide, l, t, w, h, lines):
    """Multi-line text box. lines = [(text, size, color, bold, align), ...]"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        txt = ln[0]
        sz = ln[1] if len(ln) > 1 else 14
        col = ln[2] if len(ln) > 2 else GRAY
        bld = ln[3] if len(ln) > 3 else False
        al = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = col
        p.font.bold = bld
        p.font.name = "Calibri"
        p.alignment = al
        p.space_after = Pt(3)
    return box


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    W = prs.slide_width
    H = prs.slide_height

    # ========== SLIDE 1: TITLE ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=BLUE)
    tb(s, Inches(1.2), Inches(1.6), Inches(11), Inches(0.7),
       "In-Context Learning as", sz=40, color=GRAY)
    tb(s, Inches(1.2), Inches(2.3), Inches(11), Inches(0.9),
       "Latent Space Manipulation", sz=50, color=WHITE, bold=True)
    rect(s, Inches(1.2), Inches(3.4), Inches(2.5), Pt(4), fill=CYAN)
    tb(s, Inches(1.2), Inches(3.7), Inches(11), Inches(0.6),
       "A Unified QVM-FRE-BI Framework for Understanding and Improving ICL", sz=20, color=CYAN)
    mtb(s, Inches(1.2), Inches(4.8), Inches(11), Inches(2.0), [
        ("Pratyay Dutta", 22, WHITE, True),
        ("Department of Computer Science, University of California, Riverside", 16, GRAY),
        ("Lawrence Livermore National Laboratory", 16, GRAY),
        ("", 8, GRAY),
        ("Phase 1: Systematic Problem Identification  |  June 2026", 16, BLUE),
    ])

    # ========== SLIDE 2: MOTIVATION ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=PURPLE)
    tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.6), "Motivation", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2), Pt(3), fill=PURPLE)

    # Left card
    rect(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.7), border=MGRAY)
    mtb(s, Inches(1.1), Inches(1.4), Inches(5), Inches(2.5), [
        ("The Problem", 20, BLUE, True),
        ("", 6, GRAY),
        ("  ICL works -- but we don't fully understand WHY", 14, GRAY),
        ("  Fragile: sensitive to demo order, selection, formatting", 14, GRAY),
        ("  Lacks negative feedback mechanism", 14, GRAY),
        ("  No explicit feature selection / axis weighting", 14, GRAY),
        ("  Theory-practice gap in latent space dynamics", 14, GRAY),
    ])

    # Right card
    rect(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(2.7), border=MGRAY)
    mtb(s, Inches(7.1), Inches(1.4), Inches(5), Inches(2.5), [
        ("The Opportunity", 20, GREEN, True),
        ("", 6, GRAY),
        ("  Classical CBIR solved analogous problems 20 years ago", 14, GRAY),
        ("  QVM (Rocchio) = Attention-weighted residual updates", 14, GRAY),
        ("  FRE = Per-axis feature importance weighting", 14, GRAY),
        ("  BI = Posterior task inference P(y|x,C)", 14, GRAY),
        ("  Cross-disciplinary bridge = high-impact research", 14, GRAY),
    ])

    # Research question
    rect(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.7), fill=RGBColor(0x12, 0x1A, 0x35), border=CYAN, bw=Pt(2))
    mtb(s, Inches(1.2), Inches(4.5), Inches(10.8), Inches(2.4), [
        ("Core Research Question", 22, CYAN, True),
        ("", 8, GRAY),
        ("\"Do transformer attention mechanisms during ICL implicitly perform", 17, WHITE),
        ("Feature Relevance Estimation (FRE) -- dynamically re-weighting latent", 17, WHITE),
        ("dimensions based on the variance structure of demonstrations -- and can", 17, WHITE),
        ("this be made explicit to improve ICL performance and interpretability?\"", 17, WHITE),
    ])

    # ========== SLIDE 3: BACKGROUND ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=BLUE)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Background: How ICL Works", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=BLUE)

    rect(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.0), border=MGRAY)
    mtb(s, Inches(1.1), Inches(1.4), Inches(11), Inches(1.8), [
        ("ICL Formulation (Dong et al., 2022)", 18, BLUE, True),
        ("", 4, GRAY),
        ("Parameters are frozen. Context C = {(x1,y1),...,(xk,yk)} concatenated with query x_target.", 14, GRAY),
        ("y* = argmax_y  P(y | x_target, C; theta)           -- No gradient updates, only forward pass", 16, WHITE, True),
        ("Under the hood: P(y|x,C) = integral P(y|x,gamma) P(gamma|C) d_gamma   -- Implicit Bayesian inference", 16, CYAN, True),
    ])

    # Key findings
    rect(s, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.5), border=MGRAY)
    mtb(s, Inches(1.1), Inches(3.6), Inches(5), Inches(3.3), [
        ("Key Findings from Literature", 18, WHITE, True),
        ("", 4, GRAY),
        ("  > Attention = implicit gradient descent (von Oswald '23)", 13, GRAY),
        ("  > ICL creates task vectors in middle layers (Hendel '23)", 13, GRAY),
        ("  > Function vectors in specific heads (Todd, ICLR '24)", 13, GRAY),
        ("  > Format learning can overpower semantics (Min '22)", 13, GRAY),
        ("  > ICL uses label info in certain cases (Kossen '23)", 13, GRAY),
        ("  > Context C acts as \"implicit training data\"", 13, GRAY),
    ])

    # Known problems
    rect(s, Inches(6.8), Inches(3.5), Inches(5.5), Inches(3.5), border=RED)
    mtb(s, Inches(7.1), Inches(3.6), Inches(5), Inches(3.3), [
        ("Known Failure Modes", 18, RED, True),
        ("", 4, GRAY),
        ("  x  Sensitive to demonstration ordering (recency bias)", 13, GRAY),
        ("  x  Brittle to demo selection -- small changes, big shifts", 13, GRAY),
        ("  x  No mechanism for negative demonstrations", 13, GRAY),
        ("  x  Calibration problems (overconfident wrong answers)", 13, GRAY),
        ("  x  Scales poorly with complex/many-step tasks", 13, GRAY),
        ("  x  Context rot in production deployments", 13, GRAY),
    ])

    # ========== SLIDE 4: UNIFIED FRAMEWORK ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=GREEN)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "The Unified ICL-CBIR Framework", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=GREEN)

    # QVM column
    rect(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.5), border=BLUE, bw=Pt(2))
    mtb(s, Inches(0.7), Inches(1.5), Inches(3.4), Inches(5.3), [
        ("QVM", 28, BLUE, True, PP_ALIGN.CENTER),
        ("Query Vector Movement", 13, GRAY, False, PP_ALIGN.CENTER),
        ("", 8, GRAY),
        ("Classical (Rocchio):", 14, CYAN, True),
        ("q_new = a*q_old", 13, WHITE),
        ("       + b*mean(D+)", 13, WHITE),
        ("       - g*mean(D-)", 13, WHITE),
        ("", 6, GRAY),
        ("Transformer:", 14, CYAN, True),
        ("h_new = h_old", 13, WHITE),
        ("  + Sum softmax(QK'/sqrt(d)) * V", 13, WHITE),
        ("", 8, GRAY),
        ("  [OK] Residual = a*q_old", 13, GREEN),
        ("  [OK] Attention = b*relevance", 13, GREEN),
        ("  [!!] Missing g-term (repulsion)", 13, RED),
    ])

    # FRE column
    rect(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), border=PURPLE, bw=Pt(2))
    mtb(s, Inches(4.9), Inches(1.5), Inches(3.4), Inches(5.3), [
        ("FRE", 28, PURPLE, True, PP_ALIGN.CENTER),
        ("Feature Relevance Estimation", 13, GRAY, False, PP_ALIGN.CENTER),
        ("", 8, GRAY),
        ("Classical (TPAMI 2005):", 14, CYAN, True),
        ("w_j proportional to 1/sigma_j", 13, WHITE),
        ("(inverse variance weighting)", 13, WHITE),
        ("d(q,x) = sqrt(Sum w_j*(q_j-x_j)^2)", 13, WHITE),
        ("", 6, GRAY),
        ("Transformer:", 14, CYAN, True),
        ("Multi-head attn = subspace FRE", 13, WHITE),
        ("LayerNorm = axis rescaling", 13, WHITE),
        ("", 8, GRAY),
        ("  [??] Implicit but uncontrolled", 13, ORANGE),
        ("  [??] Can we make it explicit?", 13, ORANGE),
    ])

    # BI column
    rect(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.5), border=CYAN, bw=Pt(2))
    mtb(s, Inches(9.1), Inches(1.5), Inches(3.4), Inches(5.3), [
        ("BI", 28, CYAN, True, PP_ALIGN.CENTER),
        ("Bayesian Inference", 13, GRAY, False, PP_ALIGN.CENTER),
        ("", 8, GRAY),
        ("Classical:", 14, CYAN, True),
        ("P(rel|x) ~ P(x|rel)*P(rel)", 13, WHITE),
        ("", 6, GRAY),
        ("Transformer:", 14, CYAN, True),
        ("P(y|x,C) =", 13, WHITE),
        (" integral P(y|x,g)*P(g|C) dg", 13, WHITE),
        ("Context C infers latent task g", 13, WHITE),
        ("", 8, GRAY),
        ("  [OK] Well-established analogy", 13, GREEN),
        ("  [OK] Pretraining = prior P(g)", 13, GREEN),
    ])

    # ========== SLIDE 5: FRE DEEP DIVE ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=PURPLE)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "FRE: The New Angle", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=PURPLE)

    rect(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.8), border=PURPLE)
    mtb(s, Inches(1.1), Inches(1.4), Inches(5), Inches(2.6), [
        ("Mathematical Formulation (TPAMI 2005)", 17, PURPLE, True),
        ("", 4, GRAY),
        ("Given relevant set R = {I_1, ..., I_n}:", 14, GRAY),
        ("", 3, GRAY),
        ("sigma_j = std(feature j across R)", 15, WHITE, True),
        ("w_j = (1/sigma_j) / Sum_k(1/sigma_k)", 15, WHITE, True),
        ("", 4, GRAY),
        ("Tight cluster --> high weight (discriminative)", 13, GREEN),
        ("Wide spread --> low weight (uninformative)", 13, ORANGE),
    ])

    rect(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(2.8), border=BLUE)
    mtb(s, Inches(7.1), Inches(1.4), Inches(5), Inches(2.6), [
        ("Transformer Analogues", 17, BLUE, True),
        ("", 4, GRAY),
        ("Multi-Head Attention:", 14, WHITE, True),
        ("   Each head = FRE on a subspace of d_model", 13, GRAY),
        ("", 3, GRAY),
        ("Layer Normalization:", 14, WHITE, True),
        ("   Rescales axis-wise variance (implicit FRE)", 13, GRAY),
        ("", 3, GRAY),
        ("Sparse Autoencoders (Anthropic '24):", 14, WHITE, True),
        ("   Decompose into monosemantic axes --> explicit FRE", 13, GRAY),
    ])

    rect(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.6), fill=RGBColor(0x1A, 0x14, 0x30), border=PURPLE, bw=Pt(2))
    mtb(s, Inches(1.2), Inches(4.55), Inches(10.8), Inches(2.4), [
        ("Key Insight: FRE Bridges Classical IR and Modern Interpretability", 19, PURPLE, True),
        ("", 6, GRAY),
        ("FRE (1/sigma weighting)  <-->  Attention Scores  <-->  SAE Feature Activations", 16, WHITE),
        ("", 3, GRAY),
        ("All three solve the same fundamental problem:", 15, GRAY),
        ("\"Which dimensions of the representation space are relevant for this query/task?\"", 16, CYAN, True),
        ("", 4, GRAY),
        ("If ICL performs implicit FRE well --> explain successes. If poorly --> improve it.", 15, GREEN),
    ])

    # ========== SLIDE 6: RESEARCH GAPS ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=RED)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Research Gaps Identified", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2), Pt(3), fill=RED)

    gaps = [
        ("Gap 1: Missing gamma-Term", RED,
         "Standard ICL has NO negative demonstration mechanism.\nThe repulsive component of classical RF is absent.",
         "Contrastive ICL (2024) begins addressing\nthis but lacks QVM theoretical grounding."),
        ("Gap 2: Implicit FRE is Uncontrolled", ORANGE,
         "Transformers may perform FRE implicitly via attention\nheads, but this is not verified or controllable.",
         "No study connects classical FRE weighting\n(1/sigma_j) to attention dynamics."),
        ("Gap 3: Positional Bias vs. Content", PURPLE,
         "Attention conflates positional proximity with semantic\nrelevance. Recency bias corrupts QVM weighting.",
         "QVM uses content-based weights. ICL uses\nposition-contaminated softmax scores."),
        ("Gap 4: No Unified Framework", BLUE,
         "QVM, FRE, and BI are studied in isolation. No one\nhas unified them for understanding ICL.",
         "Our contribution: the first unified\nQVM-FRE-BI lens for ICL mechanism analysis."),
    ]
    for i, (title, color, left_t, right_t) in enumerate(gaps):
        y = Inches(1.3) + i * Inches(1.45)
        rect(s, Inches(0.8), y, Inches(11.5), Inches(1.3), border=color)
        tb(s, Inches(1.0), y + Inches(0.05), Inches(4), Inches(0.3), title, sz=15, color=color, bold=True)
        tb(s, Inches(1.0), y + Inches(0.35), Inches(5), Inches(0.9), left_t, sz=12, color=GRAY)
        tb(s, Inches(6.8), y + Inches(0.35), Inches(5), Inches(0.9), right_t, sz=12, color=WHITE)

    # ========== SLIDE 7: EXPERIMENTAL PROTOCOL ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=CYAN)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Phase 1: Experimental Protocol", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=CYAN)
    tb(s, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
       "6 diagnostic modules to expose ICL failure modes through the QVM-FRE lens", sz=15, color=GRAY)

    modules = [
        ("M1", "Latent Noise\nInjection", BLUE, "P2", GREEN,
         "Gaussian noise in\nvalue vectors -->\nsensitivity map"),
        ("M2", "Attention\nSkewing", CYAN, "P2", GREEN,
         "Permute demo order\n--> recency bias\nquantification"),
        ("M3", "Label Space\nFlipping", GREEN, "P1", ORANGE,
         "Flip labels -->\nformat vs semantics\nstrength test"),
        ("M4", "FRE\nProbing", PURPLE, "P0", RED,
         "Per-dim ablation\n--> implicit FRE\nquality measure"),
        ("M5", "Gamma-Term\nExperiments", RED, "P0", RED,
         "Negative demos +\nsynthetic repulsion\ninjection"),
        ("M6", "Task Vector\nGeometry", ORANGE, "P1", ORANGE,
         "Extract & analyze\ntask vector\nstructure"),
    ]
    for i, (lbl, title, col, pri, pcol, desc) in enumerate(modules):
        x = Inches(0.4) + i * Inches(2.1)
        y = Inches(1.7)
        rect(s, x, y, Inches(1.95), Inches(3.5), border=col, bw=Pt(2))
        rect(s, x + Inches(0.1), y + Inches(0.1), Inches(0.55), Inches(0.35), fill=col)
        tb(s, x + Inches(0.1), y + Inches(0.1), Inches(0.55), Inches(0.35),
           lbl, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.1), y + Inches(0.55), Inches(1.75), Inches(0.7),
           title, sz=14, color=WHITE, bold=True)
        tb(s, x + Inches(0.1), y + Inches(1.4), Inches(1.75), Inches(1.3),
           desc, sz=12, color=GRAY)
        # Priority badge
        rect(s, x + Inches(0.6), y + Inches(2.9), Inches(0.7), Inches(0.35), fill=pcol)
        tb(s, x + Inches(0.6), y + Inches(2.9), Inches(0.7), Inches(0.35),
           pri, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    mtb(s, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), [
        ("Priority:   P0 = Most Novel (do first)    P1 = High Value    P2 = Supporting Evidence", 14, GRAY),
        ("", 4, GRAY),
        ("Convergence --> Cross-module analysis --> Unified sensitivity matrix --> Problem formulation", 15, CYAN, True),
    ])

    # ========== SLIDE 8: MODULE 4 FRE PROBING ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=PURPLE)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Module 4: FRE Probing -- The Core Experiment", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=PURPLE)

    # 4a
    rect(s, Inches(0.5), Inches(1.3), Inches(3.8), Inches(2.8), border=PURPLE)
    mtb(s, Inches(0.7), Inches(1.4), Inches(3.4), Inches(2.6), [
        ("4a: Per-Dimension Sensitivity", 15, PURPLE, True),
        ("", 4, GRAY),
        ("Method:", 13, WHITE, True),
        ("Ablate each dimension j of h_q", 12, GRAY),
        ("Measure accuracy drop = importance", 12, GRAY),
        ("", 3, GRAY),
        ("Output:", 13, WHITE, True),
        ("Ground-truth FRE weight vector", 12, GREEN),
        ("fre[j] = acc_full - acc_ablated[j]", 12, WHITE),
    ])

    # 4b
    rect(s, Inches(4.6), Inches(1.3), Inches(3.8), Inches(2.8), border=BLUE)
    mtb(s, Inches(4.8), Inches(1.4), Inches(3.4), Inches(2.6), [
        ("4b: Attention Head FRE", 15, BLUE, True),
        ("", 4, GRAY),
        ("Method:", 13, WHITE, True),
        ("Measure each head's contribution", 12, GRAY),
        ("to each dimension of output", 12, GRAY),
        ("", 3, GRAY),
        ("Output:", 13, WHITE, True),
        ("Head-FRE correlation score", 12, GREEN),
        ("pearsonr(head_contrib, fre)", 12, WHITE),
    ])

    # 4c
    rect(s, Inches(8.7), Inches(1.3), Inches(4.0), Inches(2.8), border=CYAN)
    mtb(s, Inches(8.9), Inches(1.4), Inches(3.6), Inches(2.6), [
        ("4c: SAE Feature Decomposition", 15, CYAN, True),
        ("", 4, GRAY),
        ("Method:", 13, WHITE, True),
        ("Use Sparse Autoencoders to find", 12, GRAY),
        ("monosemantic features changed by ICL", 12, GRAY),
        ("", 3, GRAY),
        ("Output:", 13, WHITE, True),
        ("ICL FRE map in interpretable space", 12, GREEN),
        ("delta = SAE(h_ICL) - SAE(h_zero)", 12, WHITE),
    ])

    # Expected outcomes
    rect(s, Inches(0.5), Inches(4.4), Inches(12.2), Inches(2.6), border=MGRAY)
    mtb(s, Inches(0.8), Inches(4.5), Inches(11.6), Inches(2.4), [
        ("Expected Outcomes --> Problem Formulation", 17, WHITE, True),
        ("", 4, GRAY),
        ("A:  Strong implicit FRE --> Importance concentrated on task-relevant axes --> ICL feature selection works", 13, GREEN),
        ("B:  Weak/noisy FRE --> Importance diffuse across all dims --> PROBLEM: ICL wastes capacity", 13, RED),
        ("C:  Unstable FRE --> Weights vary wildly across demo sets --> PROBLEM: implicit FRE not robust", 13, ORANGE),
        ("", 4, GRAY),
        ("Each outcome maps to a specific research direction for Phase 2 solution design.", 14, CYAN),
    ])

    # ========== SLIDE 9: MODULE 5 GAMMA ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=RED)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Module 5: The Missing Gamma-Term", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=RED)

    rect(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5), border=RED)
    mtb(s, Inches(1.1), Inches(1.4), Inches(5), Inches(2.3), [
        ("The Missing Component", 17, RED, True),
        ("", 4, GRAY),
        ("Rocchio QVM:", 14, CYAN, True),
        ("q_new = a*q + b*mean(D+) - g*mean(D-)", 15, WHITE, True),
        ("          |           |           |", 14, MGRAY),
        ("      residual   attraction  REPULSION", 14, GRAY),
        ("", 3, GRAY),
        ("ICL Attention:", 14, CYAN, True),
        ("h_new = h_old + Sum(attn*V)     ???", 15, WHITE, True),
        ("          |          |        NO g-TERM!", 14, RED, True),
    ])

    rect(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(2.5), border=MGRAY)
    mtb(s, Inches(7.1), Inches(1.4), Inches(5), Inches(2.3), [
        ("Experimental Design", 17, CYAN, True),
        ("", 4, GRAY),
        ("5a: Contrastive ICL", 14, WHITE, True),
        ("   Add negative demos --> measure repulsion", 13, GRAY),
        ("   in latent space via cosine similarity", 13, GRAY),
        ("", 3, GRAY),
        ("5b: Synthetic Gamma-Term Injection", 14, WHITE, True),
        ("   h_modified = h_q - gamma * neg_centroid", 13, GRAY),
        ("   Sweep gamma in {0.1, 0.3, 0.5, 0.7, 1.0}", 13, GRAY),
        ("   Measure accuracy at each gamma level", 13, GRAY),
    ])

    rect(s, Inches(0.8), Inches(4.1), Inches(11.5), Inches(3.0), fill=RGBColor(0x20, 0x12, 0x15), border=RED, bw=Pt(2))
    mtb(s, Inches(1.2), Inches(4.25), Inches(10.8), Inches(2.7), [
        ("Hypotheses", 17, RED, True),
        ("", 4, GRAY),
        ("H1: Adding negative demonstrations will cause the query vector to move AWAY from", 14, WHITE),
        ("    incorrect class centroids in latent space (measurable via cosine distance increase).", 14, WHITE),
        ("", 3, GRAY),
        ("H2: There exists an optimal gamma* that maximizes ICL accuracy. This gamma* will be", 14, WHITE),
        ("    task-dependent --> need for a LEARNED repulsion mechanism (Phase 2: RL-based selection).", 14, WHITE),
        ("", 3, GRAY),
        ("H3: The improvement from gamma-term will be larger for ambiguous/hard examples,", 14, WHITE),
        ("    where the query vector sits near the decision boundary between class manifolds.", 14, WHITE),
    ])

    # ========== SLIDE 10: EVALUATION ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=BLUE)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Evaluation Framework", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2), Pt(3), fill=BLUE)

    rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(3.0), border=MGRAY)
    mtb(s, Inches(0.7), Inches(1.4), Inches(5.6), Inches(2.8), [
        ("Benchmarks", 18, BLUE, True),
        ("", 4, GRAY),
        ("Easy:     SST-2 (2-class)  |  MR (2-class)", 14, GREEN),
        ("Medium:   AGNews (4-class) |  TREC (6-class)", 14, ORANGE),
        ("Hard:     DBPedia (14-class) | SST-5 (5-class)", 14, RED),
        ("Reason:   GSM8K subset | BoolQ", 14, PURPLE),
        ("", 3, GRAY),
        ("Standardized templates (StaICC) to isolate", 13, GRAY),
        ("mechanism from prompt engineering effects", 13, GRAY),
    ])

    rect(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(3.0), border=MGRAY)
    mtb(s, Inches(7.0), Inches(1.4), Inches(5.2), Inches(2.8), [
        ("Models (Open-Weight)", 18, BLUE, True),
        ("", 4, GRAY),
        ("GPT-2 Small (124M) -- fast iteration", 14, GRAY),
        ("GPT-2 Large (774M) -- scale comparison", 14, GRAY),
        ("Pythia-1.4B / 6.9B -- interp. support", 14, GRAY),
        ("LLaMA-2/3 (7-8B) -- production scale", 14, GRAY),
        ("", 3, GRAY),
        ("All require access to internal activations", 13, ORANGE),
        ("for probing, patching, and FRE analysis", 13, ORANGE),
    ])

    rect(s, Inches(0.5), Inches(4.6), Inches(12.2), Inches(2.4), border=MGRAY)
    mtb(s, Inches(0.7), Inches(4.7), Inches(11.8), Inches(2.2), [
        ("Metrics -- Beyond Black-Box Accuracy", 18, CYAN, True),
        ("", 4, GRAY),
        ("Performance:  Accuracy, Macro-F1, ECE (calibration)", 13, GRAY),
        ("Representation:  Cosine similarity, L2 distance, KL divergence", 13, GRAY),
        ("FRE-specific:  Dimension ablation score, FRE entropy, cross-task overlap", 13, GRAY),
        ("Geometry:  Task vector inter/intra-class distance", 13, GRAY),
        ("", 3, GRAY),
        ("All: mean +/- std over >=5 seeds  |  Paired t-tests  |  Bootstrap CIs  |  Cohen's d", 14, GREEN, True),
    ])

    # ========== SLIDE 11: TIMELINE ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=ORANGE)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Timeline & Methodology", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2), Pt(3), fill=ORANGE)

    weeks = [
        ("Week 1", "Infrastructure", "Model loading\nActivation hooks\nEval pipeline", BLUE),
        ("Week 2-3", "FRE Probing", "Per-dim sensitivity\nHead FRE analysis\nSAE decomposition", PURPLE),
        ("Week 3-4", "Gamma-Term", "Contrastive ICL\nSynthetic gamma\ninjection", RED),
        ("Week 4-5", "Perturbations", "Noise injection\nOrder permutation\nLabel flipping", CYAN),
        ("Week 5-6", "Geometry", "Task vectors\nExtraction\nStability analysis", ORANGE),
        ("Week 7-8", "Synthesis", "Cross-module\nanalysis + Problem\nformulation", GREEN),
    ]
    for i, (wk, title, desc, col) in enumerate(weeks):
        x = Inches(0.4) + i * Inches(2.1)
        rect(s, x, Inches(1.4), Inches(1.9), Inches(2.5), border=col, bw=Pt(2))
        tb(s, x + Inches(0.05), Inches(1.5), Inches(1.8), Inches(0.35),
           wk, sz=13, color=col, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.05), Inches(1.9), Inches(1.8), Inches(0.35),
           title, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.05), Inches(2.4), Inches(1.8), Inches(1.3),
           desc, sz=11, color=GRAY, align=PP_ALIGN.CENTER)

    rect(s, Inches(0.4), Inches(4.0), Inches(12.2), Pt(3), fill=ORANGE)

    # Do's and Don'ts
    rect(s, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.7), border=GREEN)
    mtb(s, Inches(0.7), Inches(4.4), Inches(5.4), Inches(2.5), [
        ("DO", 18, GREEN, True),
        ("", 3, GRAY),
        ("  Start with GPT-2, scale up to LLaMA", 13, GRAY),
        ("  Use standardized prompts (StaICC)", 13, GRAY),
        ("  Report with confidence intervals (>=5 seeds)", 13, GRAY),
        ("  Prioritize FRE (M4) and gamma-term (M5) first", 13, GRAY),
        ("  Look for cross-module patterns", 13, GRAY),
        ("  Treat null results as findings", 13, GRAY),
    ])

    rect(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.7), border=RED)
    mtb(s, Inches(7.0), Inches(4.4), Inches(5.4), Inches(2.5), [
        ("DON'T", 18, RED, True),
        ("", 3, GRAY),
        ("  Jump to solutions before diagnosis", 13, GRAY),
        ("  Overfit to one model/dataset", 13, GRAY),
        ("  Use only accuracy (need latent metrics)", 13, GRAY),
        ("  Confuse prompt engineering with mechanism", 13, GRAY),
        ("  Skip statistical rigor", 13, GRAY),
        ("  Run all experiments simultaneously", 13, GRAY),
    ])

    # ========== SLIDE 12: DECISION TREE ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=GREEN)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Problem Identification --> Solution Direction", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=GREEN)

    # Column headers
    tb(s, Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.3),
       "EXPERIMENTAL FINDING", sz=12, color=MGRAY, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Inches(4.5), Inches(1.2), Inches(3.8), Inches(0.3),
       "PROBLEM FORMULATION", sz=12, color=MGRAY, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Inches(8.8), Inches(1.2), Inches(4.0), Inches(0.3),
       "PHASE 2 DIRECTION", sz=12, color=MGRAY, bold=True, align=PP_ALIGN.CENTER)

    paths = [
        ("IF M4 shows weak FRE", "ICL lacks feature\nselection", "Inject explicit FRE via\nlearned axis weighting", PURPLE),
        ("IF M5: gamma-term helps", "ICL missing\nrepulsion mechanism", "Contrastive attention /\nrepulsive residuals", RED),
        ("IF M2: recency bias", "QVM weighting is\ncorrupted by position", "Positional debiasing /\nRocchio-inspired weights", BLUE),
        ("IF M3: format > semantics", "Format overpowers\ncontent understanding", "Layer interventions for\nsemantic features", CYAN),
        ("IF M6: unstable vectors", "Task vectors not\nrobust to demo choice", "Task vector regularization /\ndemo-agnostic encoding", ORANGE),
    ]
    for i, (cond, prob, sol, col) in enumerate(paths):
        y = Inches(1.6) + i * Inches(1.1)
        rect(s, Inches(0.5), y, Inches(3.3), Inches(0.9), border=col)
        tb(s, Inches(0.6), y + Inches(0.1), Inches(3.1), Inches(0.7), cond, sz=13, color=col, bold=True)

        rect(s, Inches(3.9), y + Inches(0.38), Inches(0.5), Pt(3), fill=col)

        rect(s, Inches(4.5), y, Inches(3.8), Inches(0.9), border=MGRAY)
        tb(s, Inches(4.6), y + Inches(0.05), Inches(3.6), Inches(0.8), prob, sz=12, color=WHITE)

        rect(s, Inches(8.4), y + Inches(0.38), Inches(0.3), Pt(3), fill=col)

        rect(s, Inches(8.8), y, Inches(4.0), Inches(0.9), border=GREEN)
        tb(s, Inches(8.9), y + Inches(0.05), Inches(3.8), Inches(0.8), sol, sz=12, color=GREEN)

    # ========== SLIDE 13: BIG PICTURE ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=CYAN)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "The Big Picture: Cross-Disciplinary Bridge", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=CYAN)

    # Left: Classical
    rect(s, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.0), border=ORANGE)
    rect(s, Inches(0.8), Inches(1.5), Inches(4.5), Inches(0.5), fill=ORANGE)
    tb(s, Inches(0.8), Inches(1.5), Inches(4.5), Inches(0.5),
       "Classical CBIR (2005)", sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    items_l = ["User feedback", "Feature variance sigma_j", "FRE weights (1/sigma_j)",
               "QVM (query movement)", "Bayesian ranking", "RL policy selection", "Long-term memory"]
    for i, t in enumerate(items_l):
        col = WHITE if i in [2, 3] else GRAY
        tb(s, Inches(1.0), Inches(2.2) + i * Inches(0.5), Inches(4.0), Inches(0.4),
           f"  >  {t}", sz=14, color=col)

    # Center
    tb(s, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.6),
       "=== formal ===", sz=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.4),
       "analogy", sz=14, color=CYAN, align=PP_ALIGN.CENTER)

    # Right: Modern
    rect(s, Inches(8.0), Inches(1.5), Inches(4.5), Inches(5.0), border=BLUE)
    rect(s, Inches(8.0), Inches(1.5), Inches(4.5), Inches(0.5), fill=BLUE)
    tb(s, Inches(8.0), Inches(1.5), Inches(4.5), Inches(0.5),
       "Modern Transformers (2024+)", sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    items_r = ["Demonstration examples", "Attention pattern stats", "Implicit per-head weighting (?)",
               "Residual stream updates", "Logit computation", "Layer/head specialization", "Pre-trained weights"]
    for i, t in enumerate(items_r):
        col = WHITE if i in [2, 3] else GRAY
        tb(s, Inches(8.2), Inches(2.2) + i * Inches(0.5), Inches(4.0), Inches(0.4),
           f"  >  {t}", sz=14, color=col)

    # Bottom
    rect(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), fill=RGBColor(0x10, 0x20, 0x30), border=CYAN)
    tb(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
       "Contribution: First systematic study viewing ICL through the unified QVM-FRE-BI lens",
       sz=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    # ========== SLIDE 14: CONTRIBUTIONS ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=GREEN)
    tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6),
       "Expected Contributions & Next Steps", sz=34, color=WHITE, bold=True)
    rect(s, Inches(0.8), Inches(1.0), Inches(2.5), Pt(3), fill=GREEN)

    rect(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.0), border=BLUE)
    mtb(s, Inches(0.8), Inches(1.4), Inches(5.2), Inches(2.8), [
        ("Phase 1 Deliverables", 19, BLUE, True),
        ("", 4, GRAY),
        ("1. Modular experimental codebase", 14, WHITE),
        ("   Activation extraction, intervention, evaluation", 12, GRAY),
        ("2. First FRE-ICL analysis", 14, WHITE),
        ("   Classical FRE <--> attention dynamics connection", 12, GRAY),
        ("3. Failure taxonomy", 14, WHITE),
        ("   Categorized ICL failures as latent space issues", 12, GRAY),
        ("4. Precise problem statement", 14, WHITE),
        ("   With quantitative evidence from 6 modules", 12, GRAY),
    ])

    rect(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(3.0), border=GREEN)
    mtb(s, Inches(7.1), Inches(1.4), Inches(5.0), Inches(2.8), [
        ("Phase 2 Directions (Tentative)", 19, GREEN, True),
        ("", 4, GRAY),
        ("A: Explicit FRE Injection", 14, PURPLE),
        ("   Learned axis weighting for ICL improvement", 12, GRAY),
        ("B: Gamma-Term Implementation", 14, RED),
        ("   Contrastive attention mechanism", 12, GRAY),
        ("C: RL-Based RF Selection", 14, ORANGE),
        ("   Mirroring IRRL (TPAMI 2005) for LLMs", 12, GRAY),
        ("D: Task Vector Regularization", 14, CYAN),
        ("   Robust, demo-agnostic task encoding", 12, GRAY),
    ])

    rect(s, Inches(0.5), Inches(4.6), Inches(12.2), Inches(2.4), fill=RGBColor(0x0F, 0x1F, 0x15), border=GREEN, bw=Pt(2))
    mtb(s, Inches(1.0), Inches(4.75), Inches(11.2), Inches(2.1), [
        ("The Vision", 22, GREEN, True),
        ("", 6, GRAY),
        ("Classical CBIR invested decades perfecting relevance feedback with QVM + FRE + BI.", 16, WHITE),
        ("Modern ICL performs these operations implicitly but incompletely.", 16, WHITE),
        ("By understanding WHERE ICL's implicit RF fails, we can design principled interventions", 16, WHITE),
        ("that bring the full power of classical relevance feedback into the transformer age.", 16, CYAN, True),
    ])

    # ========== SLIDE 15: THANK YOU ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, W, Pt(5), fill=BLUE)
    tb(s, Inches(2), Inches(2.0), Inches(9), Inches(1.0),
       "Thank You", sz=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    rect(s, Inches(5.5), Inches(3.1), Inches(2.3), Pt(4), fill=CYAN)
    tb(s, Inches(2), Inches(3.4), Inches(9), Inches(0.6),
       "Questions & Discussion", sz=28, color=GRAY, align=PP_ALIGN.CENTER)
    mtb(s, Inches(2), Inches(4.6), Inches(9), Inches(2.0), [
        ("Pratyay Dutta", 22, WHITE, True, PP_ALIGN.CENTER),
        ("pdutta@ucr.edu", 16, CYAN, False, PP_ALIGN.CENTER),
        ("", 8, GRAY),
        ("Department of Computer Science", 16, GRAY, False, PP_ALIGN.CENTER),
        ("University of California, Riverside", 16, GRAY, False, PP_ALIGN.CENTER),
        ("Lawrence Livermore National Laboratory", 16, GRAY, False, PP_ALIGN.CENTER),
    ])

    out = r"d:\pratyay\LLNL\Project\ICL\ICL_Latent_Space_Manipulation_Research.pptx"
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
