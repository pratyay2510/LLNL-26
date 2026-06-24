"""
Mentor Proposal Deck: "Steer, Then Invert"
Auditable & On-Manifold Control of In-Context Reasoning via Invertible Latent Steering.
Visual, low-clutter. Builds on the house style of create_mentor_ppt.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SW, SH = Inches(13.333), Inches(7.5)
DARK = RGBColor(0x0E, 0x12, 0x1F)
CARD = RGBColor(0x16, 0x21, 0x3E)
CARD2 = RGBColor(0x12, 0x1A, 0x35)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA8, 0xB4, 0xC8)
MG = RGBColor(0x5A, 0x68, 0x82)


def bg(s, c=DARK):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()
    r.shadow.inherit = False
    return r


def box(s, l, t, w, h, f=CARD, b=None, bw=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    r = s.shapes.add_shape(shape, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = f
    if b:
        r.line.color.rgb = b; r.line.width = bw
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def tx(s, l, t, w, h, txt, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.text = txt; p.alignment = a
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.italic = italic
    p.font.name = "Calibri"
    return bx


def ml(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln[0]; p.font.size = Pt(ln[1]); p.font.color.rgb = ln[2]
        p.font.bold = ln[3] if len(ln) > 3 else False
        p.alignment = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
        p.font.name = "Consolas" if (len(ln) > 5 and ln[5] == "mono") else "Calibri"
        p.space_after = Pt(2)
    return bx


def hdr(s, title, color, kicker=None):
    box(s, 0, 0, SW, Pt(6), f=color, shape=MSO_SHAPE.RECTANGLE)
    tx(s, Inches(0.7), Inches(0.28), Inches(12), Inches(0.7), title, 30, WHITE, True)
    box(s, Inches(0.72), Inches(0.98), Inches(2.0), Pt(3.5), f=color, shape=MSO_SHAPE.RECTANGLE)
    if kicker:
        tx(s, Inches(0.72), Inches(1.06), Inches(11.8), Inches(0.4), kicker, 13.5, GRAY)


def arrow(s, l, t, w, h, shape, c):
    a = s.shapes.add_shape(shape, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = c; a.line.fill.background()
    a.shadow.inherit = False
    return a


def main():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    L = prs.slide_layouts[6]

    # ===================== SLIDE 1: TITLE =====================
    s = prs.slides.add_slide(L); bg(s)
    box(s, 0, 0, SW, Pt(6), f=CYAN, shape=MSO_SHAPE.RECTANGLE)
    box(s, 0, SH - Pt(6), SW, Pt(6), f=PURPLE, shape=MSO_SHAPE.RECTANGLE)
    tx(s, Inches(1.1), Inches(1.5), Inches(11.1), Inches(0.7), "Steer, Then Invert", 52, WHITE, True)
    box(s, Inches(1.15), Inches(2.55), Inches(3.0), Pt(4), f=CYAN, shape=MSO_SHAPE.RECTANGLE)
    tx(s, Inches(1.1), Inches(2.75), Inches(11.2), Inches(1.0),
       "Auditable & On-Manifold Control of In-Context Reasoning\nvia Invertible Latent Steering", 23, CYAN)
    ml(s, Inches(1.1), Inches(4.5), Inches(11.2), Inches(2.0), [
        ("Pratyay Dutta", 22, WHITE, True),
        ("UC Riverside   |   Lawrence Livermore National Laboratory", 15, GRAY),
        ("", 7, GRAY),
        ("Advisor: Prof. Bir Bhanu        Mentor Meeting  ·  June 2026", 16, BLUE),
    ])

    # ===================== SLIDE 2: ONE-LINE THESIS =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "The Idea in One Line", CYAN)
    box(s, Inches(0.7), Inches(1.55), Inches(11.95), Inches(1.7), f=CARD2, b=CYAN, bw=Pt(1.75))
    ml(s, Inches(1.0), Inches(1.7), Inches(11.4), Inches(1.5), [
        ("Control a model's reasoning by moving its latent state — then INVERT the moved", 21, WHITE, True),
        ("state to read back, in plain text, exactly what prompt that reasoning corresponds to.", 21, WHITE, True),
    ], anchor=MSO_ANCHOR.MIDDLE)

    tx(s, Inches(0.7), Inches(3.55), Inches(12), Inches(0.4), "Why this matters", 18, CYAN, True)
    pts = [
        ("Not the output.", "We intervene on the reasoning state itself, not on the prompt or the final answer.", GREEN),
        ("Make it legible.", "Latent steering today is a black box — you never know what prompt your edit means.", ORANGE),
        ("Make it safe.", "We measure how far an edit pushes the model OFF the manifold of real prompts.", BLUE),
    ]
    for i, (h, d, c) in enumerate(pts):
        x = Inches(0.7) + i * Inches(4.07)
        box(s, x, Inches(4.1), Inches(3.85), Inches(2.6), f=CARD, b=c)
        box(s, x, Inches(4.1), Inches(3.85), Pt(5), f=c, shape=MSO_SHAPE.RECTANGLE)
        tx(s, x + Inches(0.25), Inches(4.35), Inches(3.4), Inches(0.5), h, 18, c, True)
        tx(s, x + Inches(0.25), Inches(4.95), Inches(3.4), Inches(1.6), d, 14.5, GRAY)

    # ===================== SLIDE 3: THREE PILLARS =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Three Results We Stand On", GREEN, "Each is strong alone — nobody has combined them")
    pillars = [
        ("In-Context Vectors", "Liu et al. · ICML 2024", BLUE,
         ["ICL = shifting latent states", "A single steering vector",
          "replaces demonstrations", "", "Gap: opaque, uncontrolled,", "may leave the manifold"]),
        ("Injectivity / SIPIT", "Nikolaou et al. · ICLR 2026", PURPLE,
         ["prompt ↔ latent is a", "BIJECTION (a.s.)", "SIPIT inverts it exactly,", "linear time", "",
          "Gap: only ever inverts", "NATURAL states"]),
        ("QVM / FRE  (Rocchio)", "Yin, Bhanu et al. · TPAMI 2005", CYAN,
         ["Move toward positives,", "AWAY from negatives (γ)", "Per-axis FRE weighting", "",
          "Gap: classical IR — never", "linked to LLM latents"]),
    ]
    for i, (t, sub, c, lines) in enumerate(pillars):
        x = Inches(0.55) + i * Inches(4.15)
        box(s, x, Inches(1.55), Inches(3.9), Inches(5.2), f=CARD, b=c, bw=Pt(1.75))
        box(s, x, Inches(1.55), Inches(3.9), Inches(1.05), f=c)
        tx(s, x + Inches(0.2), Inches(1.66), Inches(3.5), Inches(0.5), t, 19, WHITE, True, PP_ALIGN.CENTER)
        tx(s, x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(0.35), sub, 12, RGBColor(0xE6, 0xEE, 0xFF), False, PP_ALIGN.CENTER)
        ml(s, x + Inches(0.3), Inches(2.95), Inches(3.35), Inches(3.6),
           [(ln, 14.5, (GRAY if (ln.startswith("Gap") or ln in ("may leave the manifold",
            "replaces demonstrations", "NATURAL states", "linked to LLM latents")) else WHITE),
            ln.startswith("Gap")) for ln in lines])

    # ===================== SLIDE 4: THE GAP =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "The Open Problem Nobody Has Closed", RED)
    box(s, Inches(0.7), Inches(1.5), Inches(11.95), Inches(1.35), f=CARD2, b=RED)
    ml(s, Inches(1.0), Inches(1.62), Inches(11.4), Inches(1.1), [
        ("Steering moves you in latent space — but it is a one-way street.", 19, WHITE, True),
        ("You cannot say what prompt the edit means, nor whether any real prompt could produce it.", 16, GRAY),
    ], anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("Missing γ-term", "ICL attention is a convex mix over context — it can attract to positives but cannot REPEL from negatives.", RED),
        ("Unaccountable edits", "A steering vector has no known prompt. For a safety lab, an uninspectable latent edit is a non-starter.", ORANGE),
        ("Off-manifold risk", "Nothing guarantees the steered state is reachable by ANY prompt. No metric exists for 'how off-manifold'.", BLUE),
    ]
    for i, (t, d, c) in enumerate(rows):
        y = Inches(3.15) + i * Inches(1.25)
        box(s, Inches(0.7), y, Inches(11.95), Inches(1.1), f=CARD, b=c)
        box(s, Inches(0.7), y, Pt(6), Inches(1.1), f=c, shape=MSO_SHAPE.RECTANGLE)
        tx(s, Inches(1.0), y + Inches(0.13), Inches(3.2), Inches(0.85), t, 17, c, True, anchor=MSO_ANCHOR.MIDDLE)
        tx(s, Inches(4.2), y + Inches(0.13), Inches(8.2), Inches(0.85), d, 15, WHITE, anchor=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
       "Injectivity (SIPIT) solves #2 and #3 — and supplies the inverse map. That is our opening.", 14.5, CYAN, True)

    # ===================== SLIDE 5: THE LOOP (money diagram) =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Core Mechanism:  Steer → Invert → Verify", CYAN,
        "The first loop that inverts a *manipulated* hidden state")

    # four nodes in a cycle
    nb = dict(w=Inches(3.2), h=Inches(1.35))
    p_prompt = (Inches(0.9), Inches(1.85))
    p_latent = (Inches(9.25), Inches(1.85))
    p_steer = (Inches(9.25), Inches(4.55))
    p_recov = (Inches(0.9), Inches(4.55))

    def node(pos, title, sub, c):
        box(s, pos[0], pos[1], nb['w'], nb['h'], f=CARD, b=c, bw=Pt(2))
        tx(s, pos[0], pos[1] + Inches(0.18), nb['w'], Inches(0.5), title, 18, WHITE, True, PP_ALIGN.CENTER)
        tx(s, pos[0], pos[1] + Inches(0.72), nb['w'], Inches(0.5), sub, 13, c, False, PP_ALIGN.CENTER)

    node(p_prompt, "Prompt s", "demonstrations + query", BLUE)
    node(p_latent, "Latent  h", "ICL ≈ a gradient step", GREEN)
    node(p_steer, "Steered  z = h + λv", "v = +positives − γ·negatives", ORANGE)
    node(p_recov, "Recovered prompt  ŝ", "SIPIT(z) — read it aloud", PURPLE)

    # arrows around the cycle
    arrow(s, Inches(4.25), Inches(2.25), Inches(4.85), Inches(0.55), MSO_SHAPE.RIGHT_ARROW, GREEN)
    tx(s, Inches(4.25), Inches(1.78), Inches(4.85), Inches(0.4), "forward pass", 13, GREEN, True, PP_ALIGN.CENTER)
    arrow(s, Inches(10.5), Inches(3.25), Inches(0.6), Inches(1.25), MSO_SHAPE.DOWN_ARROW, ORANGE)
    tx(s, Inches(8.0), Inches(3.5), Inches(1.2), Inches(0.5), "steer", 13, ORANGE, True, PP_ALIGN.RIGHT)
    arrow(s, Inches(4.25), Inches(4.95), Inches(4.85), Inches(0.55), MSO_SHAPE.LEFT_ARROW, PURPLE)
    tx(s, Inches(4.25), Inches(5.5), Inches(4.85), Inches(0.4), "SIPIT  (exact inverse)", 13, PURPLE, True, PP_ALIGN.CENTER)
    arrow(s, Inches(2.2), Inches(3.25), Inches(0.6), Inches(1.25), MSO_SHAPE.UP_ARROW, CYAN)
    tx(s, Inches(2.9), Inches(3.5), Inches(2.0), Inches(0.5), "verify / audit", 13, CYAN, True)

    box(s, Inches(0.9), Inches(6.45), Inches(11.55), Inches(0.85), f=CARD2, b=CYAN)
    tx(s, Inches(1.1), Inches(6.55), Inches(11.2), Inches(0.7),
       "Inversion turns an opaque latent edit into a sentence you can inspect — and its residual "
       "measures how on-manifold the edit is.", 15, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)

    # ===================== SLIDE 6: METHOD 1 — CIS =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Method 1 — Contrastive In-Context Steering (the γ-term)", GREEN,
        "Realize the repulsion ICL attention is missing; ground it in QVM + 'ICL = gradient step'")
    box(s, Inches(0.7), Inches(1.7), Inches(5.85), Inches(2.4), f=CARD, b=GREEN)
    ml(s, Inches(0.95), Inches(1.85), Inches(5.4), Inches(2.2), [
        ("Classical Rocchio (QVM)", 16, GREEN, True),
        ("q_new = α q_old + β mean(D⁺) − γ mean(D⁻)", 15, WHITE, False, PP_ALIGN.LEFT, "mono"),
        ("", 6, GRAY),
        ("Our steering vector", 16, GREEN, True),
        ("v = β mean h(D⁺) − γ mean h(D⁻)", 15, WHITE, False, PP_ALIGN.LEFT, "mono"),
        ("z = h + λ v   (norm-preserved)", 15, WHITE, False, PP_ALIGN.LEFT, "mono"),
    ])
    box(s, Inches(6.75), Inches(1.7), Inches(5.9), Inches(2.4), f=CARD, b=BLUE)
    ml(s, Inches(7.0), Inches(1.85), Inches(5.45), Inches(2.2), [
        ("Why it is principled", 16, BLUE, True),
        ("•  γ = 0  recovers In-Context Vectors", 14.5, WHITE),
        ("•  γ > 0  =  the NEW repulsive term", 14.5, WHITE),
        ("•  v is a contrastive gradient step", 14.5, WHITE),
        ("   (ICL ≈ GD — von Oswald, Dai)", 13, GRAY),
        ("•  +FRE: scale axis j by 1/σⱼ (Bhanu)", 14.5, WHITE),
    ])
    box(s, Inches(0.7), Inches(4.35), Inches(11.95), Inches(2.55), f=CARD2, b=PURPLE)
    tx(s, Inches(0.95), Inches(4.5), Inches(11.4), Inches(0.4), "Three flavors of 'negative' — increasingly interesting", 16, PURPLE, True)
    negs = [
        ("Wrong-class", "repel from the\nincorrect label centroid", "classification"),
        ("Flawed reasoning", "correct CoT vs.\ncorrupted CoT", "reasoning control"),
        ("Spurious feature", "same label, different\nsurface form", "debiasing"),
    ]
    for i, (t, d, tag) in enumerate(negs):
        x = Inches(0.95) + i * Inches(3.85)
        box(s, x, Inches(5.0), Inches(3.6), Inches(1.7), f=CARD, b=MG)
        tx(s, x + Inches(0.2), Inches(5.1), Inches(3.2), Inches(0.4), t, 15.5, CYAN, True)
        tx(s, x + Inches(0.2), Inches(5.5), Inches(3.2), Inches(0.7), d, 13.5, WHITE)
        tx(s, x + Inches(0.2), Inches(6.28), Inches(3.2), Inches(0.35), tag, 12.5, ORANGE, True)

    # ===================== SLIDE 7: METHOD 2 — AUDIT + REGULARIZER =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Method 2 — Inversion-as-Audit  +  On-Manifold Loss", PURPLE,
        "Turn the inversion residual into a control metric and a training signal")
    box(s, Inches(0.7), Inches(1.7), Inches(5.85), Inches(2.55), f=CARD, b=PURPLE)
    ml(s, Inches(0.95), Inches(1.85), Inches(5.4), Inches(2.35), [
        ("Audit (read it back)", 16, PURPLE, True),
        ("ŝ = SIPIT(z)  — the prompt whose", 14.5, WHITE),
        ("natural trajectory matches z", 14.5, WHITE),
        ("", 5, GRAY),
        ("• exact SIPIT: fails off-manifold", 13.5, GRAY),
        ("• relaxed SIPIT: always returns a", 13.5, GRAY),
        ("  prompt = projection onto prompts", 13.5, GRAY),
    ])
    box(s, Inches(6.75), Inches(1.7), Inches(5.9), Inches(2.55), f=CARD, b=ORANGE)
    ml(s, Inches(7.0), Inches(1.85), Inches(5.45), Inches(2.35), [
        ("Inversion residual = on-manifold-ness", 15.5, ORANGE, True),
        ("ρ(z) = Σₜ minᵥ ‖ zₜ − hₜ(ŝ<ₜ ⊕ v) ‖", 14.5, WHITE, False, PP_ALIGN.LEFT, "mono"),
        ("", 5, GRAY),
        ("ρ ≈ 0  →  a real prompt makes z", 14, GREEN),
        ("ρ large →  fiction no prompt creates", 14, RED),
    ])
    box(s, Inches(0.7), Inches(4.5), Inches(11.95), Inches(2.4), f=CARD2, b=GREEN)
    tx(s, Inches(0.95), Inches(4.62), Inches(11.4), Inches(0.45), "The new training objective", 16, GREEN, True)
    ml(s, Inches(0.95), Inches(5.15), Inches(11.4), Inches(1.6), [
        ("L  =  L_task(z)   +   μ · ρ(z)   +   η · ‖ z − h ‖²", 20, WHITE, True, PP_ALIGN.CENTER, "mono"),
        ("", 6, GRAY),
        ("task control            stay on-manifold (NEW)            minimal edit", 14, GRAY, False, PP_ALIGN.CENTER),
    ])

    # ===================== SLIDE 8: INVERSE DESIGN =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Method 3 — Inverse Design of Prompts", ORANGE,
        "Optimize in latent space, then invert to a prompt — exact, not an RL search")
    steps = [
        ("1  Target", "pick a reasoning goal:\nmaximize margin\nm(z)=log p(y⁺|z)−log p(y⁻|z)", BLUE),
        ("2  Optimize", "ascend  m(z) − μρ(z)\nin latent space\n(on-manifold constrained)", GREEN),
        ("3  Invert", "s* = SIPIT(z*)\nthe prompt that would\ninduce that reasoning", PURPLE),
        ("4  Compare", "vs. RL example selection\n& retrieval — ours is\nexact + training-free", ORANGE),
    ]
    for i, (t, d, c) in enumerate(steps):
        x = Inches(0.6) + i * Inches(3.07)
        box(s, x, Inches(1.8), Inches(2.8), Inches(2.7), f=CARD, b=c, bw=Pt(1.75))
        box(s, x, Inches(1.8), Inches(2.8), Inches(0.6), f=c)
        tx(s, x, Inches(1.9), Inches(2.8), Inches(0.4), t, 17, WHITE, True, PP_ALIGN.CENTER)
        tx(s, x + Inches(0.2), Inches(2.65), Inches(2.45), Inches(1.7), d, 14, GRAY)
        if i < 3:
            arrow(s, x + Inches(2.82), Inches(2.95), Inches(0.22), Inches(0.4), MSO_SHAPE.RIGHT_ARROW, MG)
    box(s, Inches(0.6), Inches(4.9), Inches(12.05), Inches(1.9), f=CARD2, b=CYAN)
    ml(s, Inches(0.9), Inches(5.05), Inches(11.5), Inches(1.7), [
        ("Why it beats prior prompt discovery", 16, CYAN, True),
        ("Active Example Selection (RL, EMNLP'22) and AutoPrompt / Hard Prompts search DISCRETE prompt space "
         "and stay approximate.", 14.5, WHITE),
        ("We optimize in CONTINUOUS latent space and use SIPIT's EXACT inverse — the bijection is what makes "
         "this principled.", 14.5, WHITE),
    ])

    # ===================== SLIDE 9: EXPERIMENTS =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Experiments — The Predicted Phase Transition", BLUE,
        "Sweep steering strength λ; invert each steered state")
    # mock plot region
    box(s, Inches(0.7), Inches(1.65), Inches(6.4), Inches(4.05), f=CARD, b=MG)
    # axes
    arrow(s, Inches(1.25), Inches(1.95), Pt(2.5), Inches(3.3), MSO_SHAPE.RECTANGLE, MG)
    arrow(s, Inches(1.25), Inches(5.2), Inches(5.4), Pt(2.5), MSO_SHAPE.RECTANGLE, MG)
    tx(s, Inches(0.5), Inches(1.7), Inches(2.2), Inches(0.3), "ρ (residual)", 12, GRAY)
    tx(s, Inches(5.9), Inches(5.25), Inches(1.2), Inches(0.3), "λ →", 13, GRAY)
    # three bands
    box(s, Inches(1.35), Inches(2.05), Inches(1.75), Inches(3.1), f=RGBColor(0x10,0x2A,0x1F), shape=MSO_SHAPE.RECTANGLE)
    box(s, Inches(3.1), Inches(2.05), Inches(1.75), Inches(3.1), f=RGBColor(0x2A,0x24,0x12), shape=MSO_SHAPE.RECTANGLE)
    box(s, Inches(4.85), Inches(2.05), Inches(1.75), Inches(3.1), f=RGBColor(0x2A,0x16,0x16), shape=MSO_SHAPE.RECTANGLE)
    # a rising curve via small dots
    import math
    for k in range(22):
        xx = 1.4 + k * 0.235
        yy = 5.05 - 2.9 / (1 + math.exp(-(k - 12) * 0.7))
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xx), Inches(yy), Pt(6), Pt(6))
        d.fill.solid(); d.fill.fore_color.rgb = CYAN; d.line.fill.background(); d.shadow.inherit = False
    tx(s, Inches(1.35), Inches(5.22), Inches(1.75), Inches(0.5), "recover\noriginal", 11.5, GREEN, True, PP_ALIGN.CENTER)
    tx(s, Inches(3.1), Inches(5.22), Inches(1.75), Inches(0.5), "recover\nDIFFERENT", 11.5, ORANGE, True, PP_ALIGN.CENTER)
    tx(s, Inches(4.85), Inches(5.22), Inches(1.75), Inches(0.5), "off-\nmanifold", 11.5, RED, True, PP_ALIGN.CENTER)
    tx(s, Inches(2.55), Inches(2.1), Inches(2.0), Inches(0.4), "λ_c ≈ Δ/2", 12.5, WHITE, True)

    box(s, Inches(7.35), Inches(1.65), Inches(5.3), Inches(4.05), f=CARD, b=BLUE)
    ml(s, Inches(7.6), Inches(1.8), Inches(4.85), Inches(3.8), [
        ("Setup", 16, BLUE, True),
        ("Models  GPT-2 S/L, Pythia-1.4B,", 14, WHITE),
        ("        + Llama-3.1-8B (scale)", 14, WHITE),
        ("Tasks   SST-2, AGNews, style/", 14, WHITE),
        ("        safety transfer, mini-CoT", 14, WHITE),
        ("", 5, GRAY),
        ("Latent metrics (the point)", 15, CYAN, True),
        ("centroid movement · margin m(z)", 13.5, WHITE),
        ("residual ρ · recovery regime", 13.5, WHITE),
        ("token-edit dist of ŝ", 13.5, WHITE),
        ("", 4, GRAY),
        ("Output metrics confirm: acc, F1, ECE", 13, GRAY),
    ])
    tx(s, Inches(7.35), Inches(5.85), Inches(5.3), Inches(0.9),
       "This one figure is a clean, novel result — independent of any accuracy gain.", 13.5, GREEN, True)

    # ===================== SLIDE 10: CONTRIBUTIONS =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Contributions  &  Why It's Top-Tier", GREEN)
    cs = [
        ("C1", "First inversion of a MANIPULATED hidden state — latent edits become legible prompts.", PURPLE),
        ("C2", "ρ: a faithful on-manifold metric with a phase transition predicted by SIPIT's robustness bound.", CYAN),
        ("C3", "The γ-term (absent from ICL) improves controllable reasoning, visibly in latent space.", GREEN),
        ("C4", "On-manifold regularizer → control that is equally effective but AUDITABLE & robust.", ORANGE),
        ("C5", "Inverse design: optimize latents, invert to prompts — exact alternative to RL selection.", BLUE),
    ]
    for i, (n, d, c) in enumerate(cs):
        y = Inches(1.55) + i * Inches(0.92)
        box(s, Inches(0.7), y, Inches(11.95), Inches(0.78), f=CARD, b=c)
        box(s, Inches(0.7), y, Inches(1.0), Inches(0.78), f=c)
        tx(s, Inches(0.7), y, Inches(1.0), Inches(0.78), n, 22, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        tx(s, Inches(1.95), y, Inches(10.4), Inches(0.78), d, 15.5, WHITE, anchor=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.7),
       "Headline: the first AUDITABLE, on-manifold method to control LLM reasoning in latent space — "
       "built on the freshly-proven invertibility of transformers.  C1+C2 alone is a paper.", 14, CYAN, True)

    # ===================== SLIDE 11: PLAN =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "4-Week Sprint  →  Submission-Grade Draft", ORANGE)
    wk = [
        ("Days 1–4", "Reproduce", "ICV steer + SIPIT\ninvert in our harness\n(gate before moving on)", MG),
        ("Days 5–11", "Phase A", "Steer→Invert phase\ntransition — the SAFE\ncore result", CYAN),
        ("Days 12–20", "Phase B", "γ-term reasoning\ncontrol; latent-space\nmeasurement + audit", GREEN),
        ("Days 21–27", "Phase C", "On-manifold loss +\ninverse design\n(upside)", PURPLE),
        ("Days 28–30", "Write", "Ablations, ≥5 seeds,\nfigures, arXiv draft", ORANGE),
    ]
    for i, (d, t, x, c) in enumerate(wk):
        px = Inches(0.45) + i * Inches(2.53)
        box(s, px, Inches(1.7), Inches(2.35), Inches(3.0), f=CARD, b=c, bw=Pt(1.75))
        tx(s, px, Inches(1.82), Inches(2.35), Inches(0.35), d, 13, c, True, PP_ALIGN.CENTER)
        tx(s, px, Inches(2.2), Inches(2.35), Inches(0.4), t, 17, WHITE, True, PP_ALIGN.CENTER)
        tx(s, px + Inches(0.18), Inches(2.75), Inches(2.0), Inches(1.8), x, 13.5, GRAY, False, PP_ALIGN.CENTER)
    box(s, Inches(0.45), Inches(5.0), Inches(12.4), Inches(1.9), f=CARD2, b=BLUE)
    ml(s, Inches(0.75), Inches(5.15), Inches(11.9), Inches(1.7), [
        ("Realistic framing", 16, BLUE, True),
        ("Month 1 = a complete, submission-grade manuscript + an arXiv preprint to plant the flag "
         "(the invertibility result is brand-new — speed matters).", 14.5, WHITE),
        ("Target: ICLR 2027 (abstract ~late Sept '26); workshop + arXiv now. A 30-day camera-ready accept is "
         "NOT realistic — a strong draft is.", 14.5, GRAY),
    ])

    # ===================== SLIDE 12: RISKS =====================
    s = prs.slides.add_slide(L); bg(s)
    hdr(s, "Risks We've Already Thought Through", RED)
    rk = [
        ("Exact SIPIT won't terminate off-manifold", "Use relaxed (nearest-token) SIPIT for the audit; its failure DEFINES the off-manifold regime."),
        ("Multi-layer steering breaks inversion", "Steer + invert at the SAME single layer ℓ*; report the layer sweep."),
        ("SIPIT is slow (~28s/prompt)", "Budget on GPT-2/Pythia; subsample Llama; cache prefix forwards; gradient-guided policy."),
        ("γ gains shrink at scale", "Lead with the scale-robust mechanism (C1, C2, C4); γ helps most near decision boundaries."),
        ("'Isn't this just ICV + a known inverter?'", "No: inverting a MANIPULATED state and the on-manifold metric/loss are both new."),
    ]
    for i, (t, m) in enumerate(rk):
        y = Inches(1.55) + i * Inches(1.02)
        box(s, Inches(0.7), y, Inches(5.75), Inches(0.9), f=CARD, b=RED)
        tx(s, Inches(0.9), y, Inches(5.4), Inches(0.9), t, 14.5, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
        arrow(s, Inches(6.5), y + Inches(0.28), Inches(0.4), Inches(0.34), MSO_SHAPE.RIGHT_ARROW, GREEN)
        box(s, Inches(7.0), y, Inches(5.65), Inches(0.9), f=CARD, b=GREEN)
        tx(s, Inches(7.2), y, Inches(5.3), Inches(0.9), m, 13.5, GRAY, anchor=MSO_ANCHOR.MIDDLE)
    tx(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.5),
       "Even a negative result (steering is usually token-invisible) is a striking interpretability statement.", 14, CYAN, True)

    # ===================== SLIDE 13: THANK YOU =====================
    s = prs.slides.add_slide(L); bg(s)
    box(s, 0, 0, SW, Pt(6), f=CYAN, shape=MSO_SHAPE.RECTANGLE)
    box(s, 0, SH - Pt(6), SW, Pt(6), f=PURPLE, shape=MSO_SHAPE.RECTANGLE)
    tx(s, Inches(2), Inches(2.1), Inches(9.33), Inches(1.0), "Discussion", 50, WHITE, True, PP_ALIGN.CENTER)
    box(s, Inches(5.4), Inches(3.25), Inches(2.5), Pt(4), f=CYAN, shape=MSO_SHAPE.RECTANGLE)
    ml(s, Inches(2), Inches(3.6), Inches(9.33), Inches(1.4), [
        ("Steer the reasoning  ·  invert to read it  ·  keep it on-manifold", 19, GRAY, False, PP_ALIGN.CENTER),
    ])
    ml(s, Inches(2), Inches(4.8), Inches(9.33), Inches(1.6), [
        ("Pratyay Dutta", 22, WHITE, True, PP_ALIGN.CENTER),
        ("pdutt005@ucr.edu", 16, CYAN, False, PP_ALIGN.CENTER),
        ("", 6, GRAY),
        ("UC Riverside  |  Lawrence Livermore National Laboratory  |  Advisor: Prof. Bir Bhanu", 14, GRAY, False, PP_ALIGN.CENTER),
    ])

    out = r"c:\Work\LLNL\LLNL-26\Project\ICL\ICL_Proposal_Steer_Then_Invert.pptx"
    prs.save(out)
    print("Saved: %s (%d slides)" % (out, len(prs.slides)))


if __name__ == "__main__":
    main()
